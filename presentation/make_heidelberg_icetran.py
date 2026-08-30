"""
make_heidelberg_icetran.py
===============================================================================
Pour presentation/presentation_heidelberg.pptx (the 8-minute Heidelberg
talk, 10 slides) into the visual template of the ICETRAN 2026 deck
(poster/IcETRAN2026_final.pptx): same embedded fonts (Gothic A1 / Lovelo),
same greens, same slide furniture (title rule, logos, full-bleed closing
slide).

Content policy: EVERY word and every picture comes from
presentation_heidelberg.pptx, read programmatically - nothing is retyped,
nothing is added, the slide count stays 10, speaker notes are copied
verbatim.  Only the layout, fonts and colours change.  A check at the end
asserts that the set of paragraphs on each output slide equals the set of
paragraphs on the source slide.

Reuses the cloning / text / picture helpers of make_template_deck.py.

Run from the project root:
    python presentation/make_heidelberg_icetran.py
    -> presentation/dual_detector_talk_icetran_style.pptx
       (replaces the 14-slide deck of make_template_deck.py, which can be
        regenerated from that script at any time)
"""

import argparse
import copy
import os
import sys

from lxml import etree
from pptx import Presentation
from pptx.oxml.ns import qn
from pptx.util import Cm

sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))
import make_template_deck as tpl          # noqa: E402  (chdir's to the repo root)

SRC = os.path.join("presentation", "presentation_heidelberg.pptx")
TEMPLATE = tpl.TEMPLATE
OUT = os.path.join("presentation", "dual_detector_talk_icetran_style.pptx")
WORK = os.path.join("presentation", "figs", "heidelberg_src")

GREEN = tpl.GREEN            # 456446, the template's text green
LIGHT = tpl.LIGHT            # 9BBB9C, the template's title green
ACCENT = "FFA51F"            # the template's own highlight orange (stats slide)
WHITE = "FFFFFF"
SRC_ACCENT = "E8760C"        # the Heidelberg deck's orange -> ACCENT
SRC_MUTED = "CADCFC"         # the Heidelberg deck's pale blue on dark -> LIGHT


# --------------------------------------------------------------------------
# source reading: text runs, pictures, notes - all verbatim
# --------------------------------------------------------------------------
def src_shape(slide, name):
    for sh in slide.shapes:
        if sh.name == name:
            return sh
    raise KeyError(name)


def runs_of(sh):
    """[[(text, bold, src_color_hex_or_None), ...] per paragraph]."""
    out = []
    for p in sh.text_frame.paragraphs:
        if not p.text.strip():
            continue
        runs = []
        for r in p.runs:
            col = None
            try:
                if r.font.color is not None and r.font.color.type is not None:
                    col = str(r.font.color.rgb)
            except Exception:
                col = None
            runs.append((r.text, bool(r.font.bold), col))
        out.append(runs)
    return out


def texts_of(sh):
    """Plain paragraph strings of a source text box."""
    return ["".join(t for t, _, _ in para) for para in runs_of(sh)]


def export_pictures(prs):
    """Save every picture of the source deck; {(slide_no, n): path}."""
    os.makedirs(WORK, exist_ok=True)
    paths = {}
    for i, s in enumerate(prs.slides, 1):
        n = 0
        for sh in s.shapes:
            if sh.shape_type == 13:
                n += 1
                path = os.path.join(WORK, "s%02d_img%d.%s" % (i, n, sh.image.ext))
                with open(path, "wb") as f:
                    f.write(sh.image.blob)
                paths[(i, n)] = path
    return paths


def notes_of(slide):
    if slide.has_notes_slide:
        return slide.notes_slide.notes_text_frame.text
    return ""


# --------------------------------------------------------------------------
# multi-run paragraphs (set_paras of the template module is one run per
# paragraph; the take-home sentence needs coloured runs inside one line)
# --------------------------------------------------------------------------
def set_runs(target, paras, style, size, color, align=None, space_after=None,
             line=None):
    """paras: [[(text, {bold, color}), ...], ...]; base style = (pPr, rPr)."""
    txBody = target.text_frame._txBody
    pPr_t, rPr_t = style
    for p in txBody.findall(qn("a:p")):
        txBody.remove(p)
    for runs in paras:
        p = etree.SubElement(txBody, qn("a:p"))
        pPr = copy.deepcopy(pPr_t) if pPr_t is not None else etree.Element(qn("a:pPr"))
        p.append(pPr)
        tpl._strip_bullet(pPr)
        if align:
            pPr.set("algn", align)
        if space_after is not None:
            for el in pPr.findall(qn("a:spcAft")):
                pPr.remove(el)
            sa = etree.SubElement(pPr, qn("a:spcAft"))
            sp = etree.SubElement(sa, qn("a:spcPts"))
            sp.set("val", str(int(space_after * 100)))
        if line is not None:
            for el in pPr.findall(qn("a:lnSpc")):
                pPr.remove(el)
            ls = etree.Element(qn("a:lnSpc"))
            sp = etree.SubElement(ls, qn("a:spcPct"))
            sp.set("val", str(int(line * 1000)))
            pPr.insert(0, ls)
        for text, opt in runs:
            r = etree.SubElement(p, qn("a:r"))
            rPr = copy.deepcopy(rPr_t)
            rPr.tag = qn("a:rPr")
            r.append(rPr)
            rPr.set("sz", str(int(opt.get("size", size) * 100)))
            rPr.set("b", "true" if opt.get("bold") else "false")
            tpl._set_color(rPr, opt.get("color", color))
            t = etree.SubElement(r, qn("a:t"))
            t.text = text


def text_box(slide, x, y, w, h, paras, style, size, color=GREEN, **kw):
    tb = slide.shapes.add_textbox(Cm(x), Cm(y), Cm(w), Cm(h))
    bodyPr = tb.text_frame._txBody.find(qn("a:bodyPr"))
    for k in ("lIns", "tIns", "rIns", "bIns"):
        bodyPr.set(k, "0")
    bodyPr.set("wrap", "square")
    set_runs(tb, paras, style, size, color, **kw)
    return tb


def plain(sh, bold=None, accent_color=ACCENT, base_color=None):
    """Source text box -> run list, keeping bold and mapping accent runs."""
    paras = []
    for runs in runs_of(sh):
        line = []
        for text, b, col in runs:
            opt = {"bold": b if bold is None else bold}
            if col == SRC_ACCENT:
                opt["color"] = accent_color
            elif base_color is not None:
                opt["color"] = base_color
            line.append((text, opt))
        paras.append(line)
    return paras


# --------------------------------------------------------------------------
# slide furniture in the template look
# --------------------------------------------------------------------------
def kicker(slide, sh_src, style):
    """The small bold section label above the title (e.g. '①  THE IDEA')."""
    text_box(slide, 2.9, 0.75, 30.0, 1.6,
             plain(sh_src, bold=True, accent_color=GREEN), style,
             size=16, color=GREEN)


def headline(slide, name, sh_src, size):
    tpl.title(slide, name, texts_of(sh_src)[0], size=size, width=45.0)


def caption(slide, x, y, w, h, sh_src, style, size=15):
    text_box(slide, x, y, w, h, plain(sh_src, bold=False), style, size=size,
             color=GREEN, align="ctr", line=105)


def body(slide, x, y, w, h, sh_src, style, size=22, space_after=10):
    text_box(slide, x, y, w, h, plain(sh_src), style, size=size, color=GREEN,
             space_after=space_after, line=105)


def stat_card(slide, x, y, w, sh_src, style, value_size=28, desc_size=16):
    """A Heidelberg stat card: first paragraph = value, rest = description.
    Rendered like the template's stat columns (bold value, plain lines)."""
    paras = runs_of(sh_src)
    items = [[(t, {"bold": True, "size": value_size}) for t, _, _ in paras[0]]]
    for runs in paras[1:]:
        items.append([(t, {"bold": False, "size": desc_size}) for t, _, _ in runs])
    text_box(slide, x, y, w, 5.0, items, style, size=desc_size, color=GREEN,
             space_after=3, line=105)


def picture(slide, path, box, fit="inside", anchor="lt"):
    return tpl.set_picture(slide, tpl.shape(slide, "Freeform 2"), path,
                           box=box, fit=fit, anchor=anchor)


# --------------------------------------------------------------------------
# the deck
# --------------------------------------------------------------------------
def build(src, template, out):
    S = list(Presentation(src).slides)
    if len(S) != 10:
        sys.exit("ERROR: expected 10 source slides, found %d" % len(S))
    pics = export_pictures(Presentation(src))

    prs = Presentation(template)
    T = list(prs.slides)
    if len(T) != 9:
        sys.exit("ERROR: expected 9 template slides, found %d" % len(T))
    t_title, t_outline, t_columns, t_picture, t_cards, t_formula, \
        t_results, t_stats, t_conclusions = T

    body_style = tpl.style_of(tpl.shape(t_columns, "TextBox 4"))    # Gothic A1
    head_style = tpl.style_of(tpl.shape(t_columns, "TextBox 3"))    # Gothic A1 Bold
    white_style = tpl.style_of(tpl.shape(t_outline, "TextBox 4"))   # white on green

    def picture_slide(src_slide, title_size=48):
        s = tpl.clone_slide(prs, t_picture)
        tpl.remove(s, "Freeform 4", "Freeform 5")
        kicker(s, src_shape(src_slide, "Text 0"), head_style)
        headline(s, "TextBox 6", src_shape(src_slide, "Text 1"), title_size)
        tpl.set_notes(s, notes_of(src_slide))
        return s

    # ---- 1  title ----------------------------------------------------------
    src = S[0]
    s = tpl.clone_slide(prs, t_title)
    text_box(s, 2.9, 3.9, 45.0, 1.5,
             plain(src_shape(src, "Text 0"), bold=True, accent_color=GREEN),
             head_style, size=16, color=GREEN, align="ctr")
    sh = tpl.shape(s, "TextBox 6")
    sh.top, sh.height = Cm(5.5), Cm(7.0)
    set_runs(sh, [[(t, {"bold": True}) for t, _, _ in para]
                  for para in runs_of(src_shape(src, "Text 1"))],
             tpl.style_of(sh), 50, LIGHT, align="ctr")
    sh = tpl.shape(s, "TextBox 7")
    sh.top, sh.height = Cm(12.7), Cm(2.6)
    set_runs(sh, plain(src_shape(src, "Text 2"), bold=False), tpl.style_of(sh),
             19, GREEN, align="ctr", space_after=2, line=105)
    text_box(s, 2.9, 15.4, 45.0, 2.0, plain(src_shape(src, "Text 3"), bold=False),
             body_style, size=14, color=GREEN, align="ctr", space_after=1,
             line=105)
    text_box(s, 6.9, 24.2, 37.0, 2.4, plain(src_shape(src, "Text 4"), bold=False),
             body_style, size=18, color=GREEN, align="ctr")
    tpl.set_notes(s, notes_of(src))

    # ---- 2  the idea: text left, schematic right ---------------------------
    src = S[1]
    s = picture_slide(src)
    body(s, 2.9, 8.2, 24.6, 18.0, src_shape(src, "Text 2"), body_style, size=24,
         space_after=12)
    picture(s, pics[(2, 1)], box=(29.0, 7.9, 18.9, 17.4), anchor="rt")
    caption(s, 29.0, 25.6, 18.9, 1.4, src_shape(src, "Text 3"), body_style)

    # ---- 3  the measurement: photo left, three stat cards right ------------
    src = S[2]
    s = picture_slide(src)
    picture(s, pics[(3, 1)], box=(2.9, 7.9, 25.5, 14.0), anchor="lt")
    caption(s, 2.9, 22.0, 25.5, 1.4, src_shape(src, "Text 2"), body_style)
    for name, y in (("Text 4", 7.9), ("Text 6", 12.9), ("Text 8", 17.9)):
        stat_card(s, 30.5, y, 17.4, src_shape(src, name), body_style)
    text_box(s, 2.9, 24.0, 45.0, 3.4, plain(src_shape(src, "Text 9"), bold=False),
             body_style, size=17, color=GREEN)

    # ---- 4  hero result: wide figure, caption, three stat cards ------------
    src = S[3]
    s = picture_slide(src, title_size=44)
    picture(s, pics[(4, 1)], box=(7.4, 7.5, 36.0, 14.4), fit="width", anchor="ct")
    caption(s, 2.9, 22.0, 45.0, 1.4, src_shape(src, "Text 2"), body_style)
    for name, x in (("Text 4", 2.9), ("Text 6", 18.2), ("Text 8", 33.5)):
        stat_card(s, x, 23.7, 14.5, src_shape(src, name), body_style)

    # ---- 5  learned fusion: text + accent line left, showcase right --------
    src = S[4]
    s = picture_slide(src)
    body(s, 2.9, 8.0, 21.6, 15.6, src_shape(src, "Text 2"), body_style, size=20,
         space_after=8)
    text_box(s, 2.9, 24.2, 21.6, 2.6, plain(src_shape(src, "Text 3"), bold=True),
             head_style, size=22, color=ACCENT)
    picture(s, pics[(5, 1)], box=(25.5, 7.9, 22.4, 15.4), anchor="rt")
    caption(s, 25.5, 22.6, 22.4, 2.4, src_shape(src, "Text 4"), body_style,
            size=14)

    # ---- 6  topography: lead line, wide relief map, three stat cards -------
    src = S[5]
    s = picture_slide(src)
    body(s, 2.9, 7.8, 45.0, 2.8, src_shape(src, "Text 2"), body_style, size=21)
    picture(s, pics[(6, 1)], box=(2.9, 10.7, 45.0, 10.4), fit="width", anchor="ct")
    caption(s, 2.9, 21.1, 45.0, 1.4, src_shape(src, "Text 3"), body_style)
    for name, x in (("Text 5", 2.9), ("Text 7", 18.2), ("Text 9", 33.5)):
        stat_card(s, x, 22.9, 14.5, src_shape(src, name), body_style)

    # ---- 7  error budget: text left, sensitivity plot right ----------------
    src = S[6]
    s = picture_slide(src)
    body(s, 2.9, 8.0, 22.6, 17.5, src_shape(src, "Text 2"), body_style, size=21)
    picture(s, pics[(7, 1)], box=(26.5, 7.9, 21.4, 14.0), anchor="rt")
    caption(s, 26.5, 20.7, 21.4, 2.4, src_shape(src, "Text 3"), body_style,
            size=14)

    # ---- 8  take-home: the template's full-bleed green closing slide -------
    src = S[7]
    s = tpl.clone_slide(prs, t_outline)
    tpl.remove(s, "TextBox 4", "TextBox 5", "TextBox 6", "TextBox 7",
               "TextBox 8", "TextBox 9", "AutoShape 10", "AutoShape 11",
               "AutoShape 12", "AutoShape 13", "AutoShape 14", "AutoShape 15")
    sh = tpl.shape(s, "TextBox 3")
    sh.left, sh.top, sh.width, sh.height = Cm(2.9), Cm(2.4), Cm(45.0), Cm(2.0)
    set_runs(sh, plain(src_shape(src, "Text 0"), bold=True, accent_color=WHITE,
                       base_color=WHITE),
             white_style, 18, WHITE, align="ctr")
    text_box(s, 4.4, 5.6, 42.0, 11.5,
             plain(src_shape(src, "Text 1"), base_color=WHITE),
             white_style, size=34, color=WHITE, align="ctr", line=112)
    text_box(s, 4.4, 17.8, 42.0, 2.4,
             plain(src_shape(src, "Text 2"), bold=True, base_color=WHITE),
             white_style, size=24, color=WHITE, align="ctr")
    text_box(s, 4.4, 21.2, 42.0, 2.0,
             plain(src_shape(src, "Text 3"), bold=False, base_color=WHITE),
             white_style, size=20, color=WHITE, align="ctr")
    text_box(s, 4.4, 25.0, 42.0, 2.4,
             plain(src_shape(src, "Text 4"), bold=False, base_color=WHITE),
             white_style, size=14, color=WHITE, align="ctr")
    tpl.set_notes(s, notes_of(src))

    # ---- 9  backup: registration ------------------------------------------
    src = S[8]
    s = picture_slide(src)
    body(s, 2.9, 8.0, 25.6, 18.0, src_shape(src, "Text 2"), body_style, size=20)
    picture(s, pics[(9, 1)], box=(30.0, 7.9, 17.9, 17.4), anchor="rt")
    caption(s, 30.0, 25.6, 17.9, 1.4, src_shape(src, "Text 3"), body_style,
            size=14)

    # ---- 10  backup: per-line fusion gains ---------------------------------
    src = S[9]
    s = picture_slide(src, title_size=44)
    picture(s, pics[(10, 1)], box=(4.4, 7.6, 42.0, 15.3), fit="width", anchor="ct")
    body(s, 2.9, 23.4, 45.0, 4.8, src_shape(src, "Text 2"), body_style, size=20,
         space_after=6)

    # ---- drop the nine template slides -------------------------------------
    for t in T:
        tpl.delete_slide(prs, t)
    prs.save(out)
    print("saved %s (%d slides)" % (out, len(prs.slides)))
    return out


# --------------------------------------------------------------------------
# verification: the words on every output slide == the words on the source
# --------------------------------------------------------------------------
def verify(src, out):
    def para_set(slide):
        ps = set()
        for sh in slide.shapes:
            if sh.has_text_frame:
                for p in sh.text_frame.paragraphs:
                    t = "".join(r.text for r in p.runs).strip()
                    if t:
                        ps.add(t)
        return ps

    A = list(Presentation(src).slides)
    B = list(Presentation(out).slides)
    assert len(A) == len(B), "slide count differs: %d vs %d" % (len(A), len(B))
    ok = True
    for i, (a, b) in enumerate(zip(A, B), 1):
        pa, pb = para_set(a), para_set(b)
        na = sum(1 for sh in a.shapes if sh.shape_type == 13)
        nb = sum(1 for sh in b.shapes if sh.shape_type == 13
                 or sh._element.find(".//" + qn("a:blip")) is not None)
        missing, extra = pa - pb, pb - pa
        if missing or extra or (notes_of(a).strip() != notes_of(b).strip()):
            ok = False
            print("slide %d: MISMATCH" % i)
            for m in sorted(missing):
                print("   missing: %r" % m[:80])
            for e in sorted(extra):
                print("   extra:   %r" % e[:80])
            if notes_of(a).strip() != notes_of(b).strip():
                print("   notes differ")
        else:
            print("slide %2d: text OK (%d paragraphs), pictures %d -> %d, notes OK"
                  % (i, len(pa), na, nb))
    print("VERIFY", "PASSED" if ok else "FAILED")
    return ok


if __name__ == "__main__":
    ap = argparse.ArgumentParser(description=__doc__.split("\n")[2])
    ap.add_argument("--src", default=SRC)
    ap.add_argument("--template", default=TEMPLATE)
    ap.add_argument("--out", default=OUT)
    args = ap.parse_args()
    build(args.src, args.template, args.out)
    sys.exit(0 if verify(args.src, args.out) else 1)
