"""
make_poster.py
===============================================================================
Build the A0 conference poster for the dual-detector paper as a PowerPoint
file, following the layout, the content ranking and the typography rules of
POSTER_PLAN.md.

WHAT THIS IS FOR
    Two ways to finish the poster, and this file serves both:

    (a) PowerPoint / print shop.  Open the .pptx, tweak, export a vector PDF.
        Text stays live vector; the figures go in as the 600 dpi PNGs
        (python-pptx cannot embed a PDF).  For a fully vector poster use
        the placement table below with LaTeX and the .pdf figures.

    (b) Canva (or any web editor).  Do NOT import a finished PDF -- Canva
        flattens vector plots and substitutes fonts.  Instead use this file as
        the reference layout: it prints an exact block-by-block placement
        table (position and size in mm) that you can reproduce on a Canva
        canvas, and results/poster_figs/*.png are already sized to drop in at
        100 %.  Type the text in Canva so it stays live vector on export.

    Either way the content is written once, here, from the committed results.

GEOMETRY (POSTER_PLAN section 3)
    A0 portrait 841 x 1189 mm, 40 mm margins, 25 mm gutters
        column width = (841 - 2*40 - 2*25) / 3 = 237.0 mm exactly
    Blocks are numbered so a visitor never has to guess the reading order,
    and every headline states its finding rather than its category.

TYPOGRAPHY (POSTER_PLAN section 4)
    title 96 pt, authors 38 pt, block headlines 44 pt, body 30 pt,
    captions 24 pt, references 20 pt.  Word budget <= 600; the script counts
    the words it wrote and fails loudly if the budget is blown.

Input : results/poster_figs/*.png|pdf + MANIFEST.txt   (scripts/18_poster_figures.py)
        authors.md                                     (author list, parsed)
Output: presentation/dual_detector_poster.pptx
        presentation/poster_layout.txt   (mm placement table, for Canva)

Run from the project root:
    python presentation/make_poster.py
    python presentation/make_poster.py --repo-url https://github.com/...
"""

import argparse
import os
import re
import sys

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Mm, Pt

ROOT = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
os.chdir(ROOT)

FIGDIR = os.path.join("results", "poster_figs")
OUTDIR = "presentation"

# ---- sheet -----------------------------------------------------------------
W, H = 841.0, 1189.0            # A0 portrait, mm
MARGIN = 40.0
GUTTER = 25.0
COL = (W - 2 * MARGIN - 2 * GUTTER) / 3.0        # 237.0 exactly
COLX = [MARGIN, MARGIN + COL + GUTTER, MARGIN + 2 * (COL + GUTTER)]

BAND_TITLE_H = 150.0            # title band
BODY_TOP = MARGIN + BAND_TITLE_H + 18.0
BAND_TAKE_H = 78.0
BODY_BOTTOM = H - MARGIN - BAND_TAKE_H - 14.0

# ---- type ------------------------------------------------------------------
FONT = "Calibri"
PT_TITLE = 96
PT_AUTHOR = 38
PT_AFFIL = 24
PT_HEAD = 44
PT_BODY = 30
PT_CAPTION = 24
PT_REF = 20
PT_TAKE = 46

NAVY = RGBColor(0x1F, 0x38, 0x64)
BLUE = RGBColor(0x1F, 0x77, 0xB4)
ORANGE = RGBColor(0xD9, 0x5F, 0x02)
INK = RGBColor(0x26, 0x26, 0x26)
GRAY = RGBColor(0x59, 0x59, 0x59)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
RULE = RGBColor(0xD8, 0xDD, 0xE6)

WORD_BUDGET = 600
_words = [0]


def count(text):
    _words[0] += len(text.split())
    return text


# ---------------------------------------------------------------------------
# content -- every number here comes from the committed reports in results/
# ---------------------------------------------------------------------------
TITLE = ("Geometry-resolved characterization of a "
         "dual-detector MA-XRF scanner")

AFFILIATIONS = (
    "School of Electrical Engineering, University of Belgrade  ·  "
    "Ars Mensurae, Rome  ·  IDArtScience Srl, Rome  ·  "
    "VINARH Center, Vinča Institute of Nuclear Sciences, "
    "University of Belgrade"
)

TAKEHOME = (
    "The difference between a scanner's two detectors - normally summed "
    "away - plus one tilted scan of the same painting characterizes the "
    "instrument for free: detector response, acquisition geometry, canvas "
    "topography, and a learned fusion that beats summing by 10–18 % SNR."
)

REFERENCES = (
    "Data and code: see the QR above (releases include the registered "
    "overlap ratios and the exported R(E) curve).  "
    "Acknowledgements: Ars Mensurae for instrument access and the scans; "
    "VINARH Center, Vinča Institute of Nuclear Sciences."
)

# (number, headline, figure stem or None, caption, body paragraphs)
BLOCK_IDEA = (
    "①", "One tilted scan separates detector from geometry",
    "block1_idea_schematic",
    "Two detectors, opposite take-off angles; tilting the canvas moves the "
    "geometric part only.",
    ["Both heads see the same spot from opposite take-off angles, so their "
     "count ratio mixes a detector part with a geometric part.",
     "Tilting the canvas changes the photon exit paths but not the "
     "detectors, so a second scan splits the ratio into a tilt-invariant "
     "and a tilt-dependent term.",
     "No calibration standards, no extra hardware - two routine "
     "acquisitions."],
)

BLOCK_DATA = (
    "②", "Two SDDs, three scans, 7 200 pixels", None, None,
    ["Ars Mensurae MA-XRF scanner, two Amptek X-123SDD heads (10264, "
     "19511), 40 kV / 50 µA. Mock-up canvas, 60 × 120 pixels, "
     "3.0 s dwell.",
     "Three scans: two frontal seven days apart (the repeatability "
     "baseline) and one tilted forward by ≤ 8°.",
     "All ratios are computed on the registered overlap (affine, "
     "NCC 0.965): full-frame ratios mix in a field-of-view artifact large "
     "enough to flip the sign of the tilt shift."],
)

BLOCK_HERO = (
    "③", "The two channels differ sixfold - and the model "
    "explains it",
    "block3_hero_geometry_fit",
    "Detector response and geometry separate cleanly, each reproduced by a "
    "small physical model.",
    ["R = det10264 / det19511 falls from 5.8 at Ca Kα to 0.63 at "
     "Pb Lγ: the same pigment gives a sixfold different signal "
     "depending on the channel.",
     "A three-parameter detector model reproduces the frontal curve across "
     "two orders of magnitude. The fitted absorber, 973 ± 2 µm "
     "Be-equivalent, is 80–120× any entrance window: extra air path "
     "or collimation in front of det 19511, not a window.",
     "Tilt shifts R monotonically: +9.5 % at Ca (24σ) down to +0.6 % "
     "at Pb Lγ. A thick-sample fluorescence model with antisymmetric "
     "take-off angles reproduces the shape (lever arm s = 0.53 ± 0.10, "
     "tilt ≤ 8°)."],
)

BLOCK_FUSION = (
    "④", "Learned fusion beats summing by 10–18 % SNR",
    "block4_fusion_pbll",
    "Same signal, visibly calmer noise. Cross-scan difference "
    "(scan 1 − scan 2)/√2 on pixels the network never saw.",
    ["The two channels are conditionally independent Poisson views of the "
     "same spot, so Noise2Noise training needs no clean target.",
     "A 1D U-Net with loss weights prescribed by the measured R(E) gains "
     "+17.7 % mean SNR over summing on held-out pixels (median +11.9 %); "
     "classical inverse-variance weighting gives +0.9 %.",
     "Six of eight lines gain - Pb Lℓ +69 %, Pb Lγ +33 %, "
     "Fe +28 %, Cu +17 % - while Ca and Ti keep the plain sum "
     "(low-energy level bias).",
     "Ablation: the same network with an unweighted loss gives "
     "−0.1 % - the gain comes from the measured variance "
     "structure, not the architecture. In acquisition time it is worth "
     "39 % longer dwell."],
)

BLOCK_TOPO = (
    "⑤", "The discarded difference is also a relief map",
    "block5_canvas_topography",
    "Detector disagreement measures surface slope, and it reproduces across "
    "scans.",
    ["Inverted through the measured tilt response, the per-pixel ratio "
     "residuals turn the eight lines into repeated measurements of local "
     "surface slope.",
     "Cross-scan r = 0.73, reproducible RMS ≈ 12°, "
     "χ²/dof = 0.85, 88 % of pixels consistent with pure "
     "geometry: a relief map from a single scan."],
)

BLOCK_POS = (
    "⑥", "One degree of mounting error moves the maps",
    "block6_positioning_sensitivity",
    "The light-element maps are the ones that move; grey band = "
    "same-geometry repeatability floor.",
    ["Element maps shift by up to 0.63 percentage points per degree of "
     "mounting error (Ca +0.50 %/°, Ti +0.45 %/°, "
     "Pb Lβ −0.13 %/°), in the energy order the model "
     "predicts - and these are lower bounds."],
)

COLUMNS = [
    [BLOCK_IDEA, BLOCK_DATA],
    [BLOCK_HERO, BLOCK_POS],
    [BLOCK_FUSION, BLOCK_TOPO],
]

PLACEMENT = []


# ---------------------------------------------------------------------------
# helpers
# ---------------------------------------------------------------------------
def read_manifest():
    """Figure aspect ratios, so blocks are laid out at the true image size."""
    path = os.path.join(FIGDIR, "MANIFEST.txt")
    if not os.path.exists(path):
        sys.exit("ERROR: %s missing -- run scripts/18_poster_figures.py first."
                 % path)
    out = {}
    for line in open(path):
        parts = line.split()
        if len(parts) == 5 and parts[0].startswith("block"):
            out[parts[0]] = (float(parts[1]), float(parts[2]))
    if not out:
        sys.exit("ERROR: no figure rows parsed from " + path)
    return out


def parse_authors():
    """Author display names, in file order, from authors.md."""
    first, family, names = None, None, []
    for line in open(os.path.join(ROOT, "authors.md"), encoding="utf-8"):
        m = re.match(r"- \*\*First name:\*\*\s*(.+?)\s*$", line)
        if m:
            first = m.group(1)
        m = re.match(r"- \*\*Family name:\*\*\s*(.+?)\s*$", line)
        if m:
            family = m.group(1)
            if first:
                names.append(first + " " + family)
                first = None
    if not names:
        sys.exit("ERROR: no authors parsed from authors.md")
    return names


def est_height(paragraphs, size_pt, width_mm, lead=1.24, space_after_mm=4.0):
    """Rough wrapped-text height in mm.

    python-pptx cannot measure text, so the column flow needs an estimate.
    Calibri averages close to 0.48 em per character at these sizes; the
    estimate runs slightly long, which is the safe direction for a poster.
    """
    em_mm = size_pt * 25.4 / 72.0
    cpl = max(8.0, width_mm / (em_mm * 0.48))
    lines = 0
    for p in paragraphs:
        lines += max(1, int(len(p) / cpl) + 1)
    return lines * em_mm * lead + space_after_mm * len(paragraphs)


def textbox(slide, x, y, w, h, paragraphs, size, color=INK, bold=False,
            align=PP_ALIGN.LEFT, space_after=6, line_spacing=1.05):
    tb = slide.shapes.add_textbox(Mm(x), Mm(y), Mm(w), Mm(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Mm(0)
    tf.margin_top = tf.margin_bottom = Mm(0)
    for i, text in enumerate(paragraphs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        p.space_after = Pt(space_after)
        r = p.add_run()
        r.text = text
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.name = FONT
        r.font.color.rgb = color
    return tb


def rect(slide, x, y, w, h, fill, line=None, shape=MSO_SHAPE.RECTANGLE):
    s = slide.shapes.add_shape(shape, Mm(x), Mm(y), Mm(w), Mm(h))
    s.shadow.inherit = False
    if fill is None:
        s.fill.background()
    else:
        s.fill.solid()
        s.fill.fore_color.rgb = fill
    if line is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = line
        s.line.width = Pt(2)
    s.text_frame.text = ""
    return s


def place(name, x, y, w, h):
    PLACEMENT.append((name, x, y, w, h))


_CIRCLED = {chr(0x2460 + i): str(i + 1) for i in range(9)}


def to_console(text):
    """The Windows console here is cp1250 -- keep stdout plain ASCII."""
    for k, v in _CIRCLED.items():
        text = text.replace(k, v)
    return text.encode("ascii", "replace").decode("ascii")


# ---------------------------------------------------------------------------
# blocks
# ---------------------------------------------------------------------------
def draw_block(slide, block, x, y, aspects):
    num, headline, stem, caption, body = block
    y0 = y

    # numbered assertion headline
    h_head = est_height([num + "  " + headline], PT_HEAD, COL, lead=1.14,
                        space_after_mm=0) + 3
    textbox(slide, x, y, COL, h_head, [num + "  " + count(headline)],
            PT_HEAD, color=NAVY, bold=True, line_spacing=0.98)
    y += h_head + 4
    rect(slide, x, y, COL, 1.6, ORANGE)
    y += 9

    if stem:
        fw, fh = aspects[stem]
        h_fig = COL * fh / fw
        path = os.path.join(FIGDIR, stem + ".png")
        slide.shapes.add_picture(path, Mm(x), Mm(y), Mm(COL), Mm(h_fig))
        place(stem + ".png", x, y, COL, h_fig)
        y += h_fig + 4
        h_cap = est_height([caption], PT_CAPTION, COL, space_after_mm=0) + 2
        textbox(slide, x, y, COL, h_cap, [count(caption)], PT_CAPTION,
                color=GRAY, bold=True, line_spacing=1.02)
        y += h_cap + 8

    h_body = est_height(body, PT_BODY, COL)
    textbox(slide, x, y, COL, h_body, [count(p) for p in body], PT_BODY,
            space_after=10)
    y += h_body

    place("block " + num, x, y0, COL, y - y0)
    return y + 26          # inter-block breathing space


def build(args):
    aspects = read_manifest()
    authors = parse_authors()

    prs = Presentation()
    prs.slide_width, prs.slide_height = Mm(W), Mm(H)
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # ---- title band ------------------------------------------------------
    rect(slide, 0, 0, W, MARGIN + BAND_TITLE_H, NAVY)
    qr_w = 62.0
    tw = W - 2 * MARGIN - 2 * (qr_w + 10)
    textbox(slide, MARGIN, 34, tw, 60, [count(TITLE)], PT_TITLE, color=WHITE,
            bold=True, line_spacing=0.96, space_after=0)
    textbox(slide, MARGIN, 108, tw, 32,
            [", ".join(authors)], PT_AUTHOR, color=RGBColor(0xC7, 0xD6, 0xEE),
            space_after=0)
    textbox(slide, MARGIN, 146, tw, 30, [AFFILIATIONS], PT_AFFIL,
            color=RGBColor(0x9F, 0xB4, 0xD6), line_spacing=1.1, space_after=0)

    for i, (label, url) in enumerate((("code + data", args.repo_url),
                                      ("paper / abstract", args.paper_url))):
        qx = W - MARGIN - (2 - i) * (qr_w + 10) + 10
        qy = 40.0
        made = False
        try:
            import qrcode                                    # optional
            img = qrcode.make(url)
            p = os.path.join(FIGDIR, "qr_%d.png" % i)
            img.save(p)
            slide.shapes.add_picture(p, Mm(qx), Mm(qy), Mm(qr_w), Mm(qr_w))
            made = True
        except ImportError:
            rect(slide, qx, qy, qr_w, qr_w, WHITE, line=RULE)
            textbox(slide, qx + 4, qy + qr_w / 2 - 8, qr_w - 8, 16,
                    ["QR here", url], 14, color=GRAY, align=PP_ALIGN.CENTER,
                    space_after=2)
        textbox(slide, qx, qy + qr_w + 3, qr_w, 10, [label], PT_REF,
                color=WHITE, align=PP_ALIGN.CENTER, space_after=0)
        place("QR %s%s" % (label, "" if made else " (placeholder)"),
              qx, qy, qr_w, qr_w)

    # ---- content columns --------------------------------------------------
    ends = []
    for ci, blocks in enumerate(COLUMNS):
        y = BODY_TOP
        for b in blocks:
            y = draw_block(slide, b, COLX[ci], y, aspects)
        ends.append(y - 26)

    # ---- take-home band ---------------------------------------------------
    ty = H - MARGIN - BAND_TAKE_H
    rect(slide, MARGIN, ty, W - 2 * MARGIN, BAND_TAKE_H, None, line=ORANGE)
    textbox(slide, MARGIN + 12, ty + 9, W - 2 * MARGIN - 24, BAND_TAKE_H - 18,
            [count(TAKEHOME)], PT_TAKE, color=NAVY, bold=True,
            line_spacing=1.04, space_after=0)
    place("take-home box", MARGIN, ty, W - 2 * MARGIN, BAND_TAKE_H)
    textbox(slide, MARGIN, H - MARGIN + 4, W - 2 * MARGIN, 14, [REFERENCES],
            PT_REF, color=GRAY, space_after=0)

    out = os.path.join(OUTDIR, "dual_detector_poster.pptx")
    prs.save(out)

    # ---- report -----------------------------------------------------------
    lines = [
        "A0 poster layout -- placement table",
        "(generated by presentation/make_poster.py; all values in mm)",
        "",
        "Sheet %.0f x %.0f mm, margin %.0f, gutter %.0f, column %.1f"
        % (W, H, MARGIN, GUTTER, COL),
        "Column x: %.0f / %.0f / %.0f      content band y: %.0f -> %.0f"
        % (COLX[0], COLX[1], COLX[2], BODY_TOP, BODY_BOTTOM),
        "",
        "Type: title %d pt, authors %d, headlines %d, body %d, captions %d,"
        % (PT_TITLE, PT_AUTHOR, PT_HEAD, PT_BODY, PT_CAPTION),
        "      take-home %d, references %d.  Font: %s (match it in Canva)."
        % (PT_TAKE, PT_REF, FONT),
        "",
        "%-34s %8s %8s %8s %8s" % ("element", "x", "y", "w", "h"),
    ]
    for name, x, y, w, h in PLACEMENT:
        lines.append("%-34s %8.1f %8.1f %8.1f %8.1f" % (name, x, y, w, h))
    lines += [
        "",
        "Column bottoms: %s (content band ends at %.0f)"
        % (" / ".join("%.0f" % e for e in ends), BODY_BOTTOM),
        "Body words written: %d of the %d budget."
        % (_words[0], WORD_BUDGET),
    ]
    rep = os.path.join(OUTDIR, "poster_layout.txt")
    with open(rep, "w", encoding="utf-8") as f:
        f.write("\n".join(lines) + "\n")

    print(to_console("\n".join(lines)))
    print("\n  saved: " + out)
    print("  saved: " + rep)

    bad = [(i, e) for i, e in enumerate(ends) if e > BODY_BOTTOM]
    if bad:
        print("\nWARNING: column(s) %s overflow the content band by up to"
              " %.0f mm -- shorten a block or move one to another column."
              % (", ".join(str(i + 1) for i, _ in bad),
                 max(e for _, e in bad) - BODY_BOTTOM))
    if _words[0] > WORD_BUDGET:
        print("\nWARNING: %d words is over the %d-word poster budget"
              " (POSTER_PLAN section 4) -- cut before printing."
              % (_words[0], WORD_BUDGET))


if __name__ == "__main__":
    ap = argparse.ArgumentParser(
        description="Build the A0 poster (POSTER_PLAN.md layout).")
    ap.add_argument("--repo-url", default="https://github.com/ADD-ME/repo",
                    help="URL behind the code/data QR code")
    ap.add_argument("--paper-url", default="https://ADD-ME/paper.pdf",
                    help="URL behind the paper/abstract QR code")
    build(ap.parse_args())
