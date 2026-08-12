"""
make_slides.py
Build the 8-minute conference talk (Heidelberg) as a PowerPoint file.

Slides follow the submitted abstract; numbers come from
results/detector_diff/ (scripts 06, 07, 07b) and PLAN.md Table 1.
Speaker notes with per-slide timing are embedded in the file
(View -> Notes / presenter view).

Run from the project root:
    python presentation/make_slides.py
Output:
    presentation/figs/geometry_schematic.png
    presentation/dual_detector_talk.pptx
"""

import os

import matplotlib.pyplot as plt
import numpy as np
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu, Inches, Pt

OUTDIR  = "presentation"
FIGDIR  = os.path.join(OUTDIR, "figs")
RESULTS = os.path.join("results", "detector_diff")

NAVY   = RGBColor(0x1F, 0x38, 0x64)
BLUE   = RGBColor(0x1F, 0x77, 0xB4)
ORANGE = RGBColor(0xD9, 0x5F, 0x02)
INK    = RGBColor(0x26, 0x26, 0x26)
GRAY   = RGBColor(0x59, 0x59, 0x59)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
FONT   = "Calibri"

SW, SH = Inches(13.333), Inches(7.5)


# ---------------------------------------------------------------------------
# geometry schematic (side view) for the idea slide
# ---------------------------------------------------------------------------
def make_schematic(path):
    fig, ax = plt.subplots(figsize=(6.2, 4.6))
    c_blue, c_orange, c_gray = "#1f77b4", "#d95f02", "#555555"

    # canvas, frontal (solid) and tilted forward (dashed, exaggerated,
    # rotated about its centre)
    ax.plot([0, 0], [-1.6, 1.6], color=c_gray, lw=5,
            solid_capstyle="round")
    th = np.radians(14)
    ex, ey = 1.6 * np.sin(th), 1.6 * np.cos(th)
    ax.plot([ex, -ex], [-ey, ey], color=c_orange, lw=3, ls="--")
    ax.text(0.14, 1.42, "canvas", fontsize=11, color=c_gray)
    ax.text(0.52, -1.62, "tilted (~8\N{DEGREE SIGN})", fontsize=11,
            color=c_orange, ha="left")
    ax.text(-0.52, 1.74, "\N{GREEK SMALL LETTER THETA}", fontsize=13,
            color=c_orange)

    # X-ray beam along the surface normal
    ax.annotate("", xy=(-0.06, 0), xytext=(-3.4, 0),
                arrowprops=dict(arrowstyle="-|>", color="#333333", lw=2))
    ax.text(-3.35, 0.13, "X-ray beam", fontsize=11, color="#333333")

    # two SDDs, symmetric about the beam
    for sgn, name in ((1, "SDD 10264"), (-1, "SDD 19511")):
        dx = -2.3 * np.cos(np.radians(45))
        dy = sgn * 2.3 * np.sin(np.radians(45))
        ax.annotate("", xy=(dx, dy), xytext=(0, 0),
                    arrowprops=dict(arrowstyle="-|>", color=c_blue,
                                    lw=1.8))
        ax.add_patch(plt.Rectangle((dx - 0.34, dy - 0.17 + sgn * 0.17),
                                   0.68, 0.34, facecolor=c_blue,
                                   edgecolor="none"))
        ax.text(dx, dy + sgn * 0.56, name, fontsize=10.5, color=c_blue,
                ha="center", va="center")
        ax.text(0.55 * dx, 0.62 * dy,
                "\N{GREEK SMALL LETTER PSI}\N{SUBSCRIPT ONE}" if sgn > 0
                else "\N{GREEK SMALL LETTER PSI}\N{SUBSCRIPT TWO}",
                fontsize=12, color=c_blue)

    ax.set_xlim(-3.8, 1.6)
    ax.set_ylim(-2.5, 2.5)
    ax.set_aspect("equal")
    ax.axis("off")
    fig.tight_layout()
    fig.savefig(path, dpi=300, bbox_inches="tight")
    plt.close(fig)


# ---------------------------------------------------------------------------
# pptx helpers
# ---------------------------------------------------------------------------
def blank_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def add_text(slide, left, top, width, height, runs_list, align=None,
             space_after=6):
    """runs_list: list of paragraphs; each is a list of
    (text, size_pt, bold, color, italic) run tuples, or the string
    "GAP" for a half-height spacer."""
    box = slide.shapes.add_textbox(left, top, width, height)
    tf  = box.text_frame
    tf.word_wrap = True
    first = True
    for para in runs_list:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.space_after = Pt(space_after)
        if align is not None:
            p.alignment = align
        if para == "GAP":
            r = p.add_run()
            r.text = " "
            r.font.size = Pt(8)
            continue
        for text, size, bold, color, italic in para:
            r = p.add_run()
            r.text = text
            r.font.name = FONT
            r.font.size = Pt(size)
            r.font.bold = bold
            r.font.italic = italic
            r.font.color.rgb = color
    return box


def add_title(slide, text, sub=None):
    runs = [[(text, 30, True, NAVY, False)]]
    if sub:
        runs.append([(sub, 16, False, GRAY, True)])
    add_text(slide, Inches(0.55), Inches(0.28), SW - Inches(1.1),
             Inches(1.0), runs, space_after=2)


def bullets(items):
    """items: list of (level, text_or_runs). Returns runs_list."""
    out = []
    for level, item in items:
        if item == "GAP":
            out.append("GAP")
            continue
        mark = "•  " if level == 0 else "     –  "
        size = 17 if level == 0 else 15
        runs = [(mark, size, False, GRAY, False)]
        if isinstance(item, str):
            runs.append((item, size, False, INK, False))
        else:
            runs += item
        out.append(runs)
    return out


def add_notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text


# ---------------------------------------------------------------------------
# build
# ---------------------------------------------------------------------------
if __name__ == "__main__":
    os.makedirs(FIGDIR, exist_ok=True)
    schematic = os.path.join(FIGDIR, "geometry_schematic.png")
    make_schematic(schematic)

    prs = Presentation()
    prs.slide_width, prs.slide_height = SW, SH

    # ---- 1 title ------------------------------------------------------
    s = blank_slide(prs)
    add_text(s, Inches(0.9), Inches(1.5), SW - Inches(1.8), Inches(1.7),
             [[("Geometry-resolved Characterization of a "
                "Dual-detector MA-XRF Scanner", 34, True, NAVY, False)]],
             align=PP_ALIGN.CENTER)
    add_text(s, Inches(0.9), Inches(3.5), SW - Inches(1.8), Inches(2.6),
             [[("Dimitrije Pešić", 18, True, INK, False),
               (",  Janko Vukobratović,  Aleksandra Stojanović,"
                "  Giulia Ristori,  Stefano Ridolfi,"
                "  Maja Gajić-Kvaščev,  Goran Kvaš"
                "čev", 18, False, INK, False)],
              "GAP",
              [("University of Belgrade - School of Electrical "
                "Engineering   ·   Ars Mensurae, Rome   ·   "
                "IDArtScience, Rome   ·   Vinča Institute of "
                "Nuclear Sciences / VINARH", 13, False, GRAY, False)],
              "GAP",
              [("EuCAIFCon 2026 - Heidelberg", 14, False, GRAY, True)]],
             align=PP_ALIGN.CENTER)
    add_notes(s, "0:00-0:15\n"
              "Good morning. I will show how a dual-detector MA-XRF "
              "scanner can be characterized - efficiencies, geometry, "
              "even the mounting angle - using nothing but routine "
              "scans of a painting.")

    # ---- 2 motivation -------------------------------------------------
    s = blank_slide(prs)
    add_title(s, "Two detectors, one discarded signal")
    add_text(s, Inches(0.9), Inches(1.55), SW - Inches(1.8), Inches(3.2),
             bullets([
                 (0, "MA-XRF scans a painting pixel by pixel; XRF lines "
                     "give element maps (Ca, Fe, Pb, ...)"),
                 (0, "Many scanners carry two SDD detectors; their "
                     "signals are summed for signal-to-noise"),
                 (0, "The difference between the two channels is "
                     "treated as noise and thrown away"),
             ]), space_after=14)
    add_text(s, Inches(0.9), Inches(4.6), SW - Inches(1.8), Inches(1.2),
             [[("That discarded difference is a calibration signal.",
                26, True, ORANGE, False)]], align=PP_ALIGN.CENTER)
    add_notes(s, "0:15-1:15\n"
              "Set the scene: MA-XRF, two detectors, summing is "
              "standard practice. The premise of this work: the part "
              "everyone throws away carries quantitative information "
              "about the instrument itself. Everything that follows is "
              "extracted from that difference.")

    # ---- 3 idea -------------------------------------------------------
    s = blank_slide(prs)
    add_title(s, "One extra tilted scan splits the physics")
    s.shapes.add_picture(schematic, Inches(0.55), Inches(1.7),
                         height=Inches(4.6))
    add_text(s, Inches(7.0), Inches(1.8), Inches(5.7), Inches(4.8),
             bullets([
                 (0, "Scan the same painting twice: frontal, and with "
                     "the canvas tilted forward"),
                 (0, "Tilting changes the photon paths to each "
                     "detector - the detectors stay the same"),
                 (0, [("The per-element ratio R(E) = det1 / det2 then "
                       "separates into:", 17, False, INK, False)]),
                 (1, [("detector part", 15, True, BLUE, False),
                      (" - tilt-invariant (absorbers, Si thickness, "
                       "gain)", 15, False, INK, False)]),
                 (1, [("geometric part", 15, True, ORANGE, False),
                      (" - tilt-dependent, energy-structured",
                       15, False, INK, False)]),
                 (0, "GAP"),
                 (0, [("“with no dedicated calibration "
                       "measurements”", 17, True, ORANGE, True)]),
             ]), space_after=10)
    add_notes(s, "1:15-2:15\n"
              "The one-slide method. The tilt is the control knob: it "
              "moves only the geometry. Comparing tilted vs frontal "
              "ratios cleanly separates detector properties from "
              "acquisition geometry. Quote the abstract: no dedicated "
              "calibration measurements - no reference targets, no "
              "monochromatic sources.")

    # ---- 4 instrument & data -----------------------------------------
    s = blank_slide(prs)
    add_title(s, "Instrument and data")
    add_text(s, Inches(0.9), Inches(1.5), SW - Inches(1.8), Inches(2.0),
             bullets([
                 (0, "Canvas copy of the Creation of Adam; Ars Mensurae "
                     "MA-XRF scanner, two SDDs (s/n 10264, 19511)"),
                 (0, "Three scans: two frontal, 7 days apart "
                     "(repeatability baseline), one tilted forward"),
                 (0, "8 usable lines from Ca Kα (3.7 keV) to "
                     "Pb Lγ (14.8 keV); bootstrap uncertainties"),
             ]), space_after=8)
    s.shapes.add_picture(
        os.path.join("results", "10264", "prova1", "element_maps.png"),
        Inches(0.65), Inches(3.55), width=Inches(12.0))
    add_notes(s, "2:15-3:00\n"
              "The test object is a canvas copy of the Creation of "
              "Adam - the element maps show the two hands. Three "
              "routine scans, nothing else. The two frontal scans give "
              "the repeatability floor; the tilted one is the "
              "geometry probe.")

    # ---- 5 table ------------------------------------------------------
    s = blank_slide(prs)
    add_title(s, "The two channels disagree - up to sixfold")
    # Ratios restricted to the registered overlap region (script 08):
    # the full-frame values mixed in a field-of-view artifact.
    rows = [
        ("Ca Kα", "3.69", "5.78", "+9.5 %", "24σ"),
        ("Ti Kα", "4.51", "2.42", "+7.9 %", "13.5σ"),
        ("Fe Kα", "6.40", "1.21", "+3.5 %", "9.2σ"),
        ("Cu Kα", "8.04", "1.01", "+2.5 %", "4.6σ"),
        ("Pb Ll", "9.19", "0.828", "+4.1 %", "5.6σ"),
        ("Pb Lα", "10.54", "0.765", "+2.2 %", "11.5σ"),
        ("Pb Lβ", "12.61", "0.699", "+1.8 %", "10.3σ"),
        ("Pb Lγ", "14.77", "0.630", "+0.6 %", "1.8σ"),
    ]
    tbl = s.shapes.add_table(9, 5, Inches(0.7), Inches(1.6),
                             Inches(6.6), Inches(4.9)).table
    heads = ("line", "E (keV)", "R frontal", "tilt shift", "signif.")
    for j, h in enumerate(heads):
        c = tbl.cell(0, j)
        c.text = h
        c.fill.solid()
        c.fill.fore_color.rgb = NAVY
        r = c.text_frame.paragraphs[0].runs[0]
        r.font.name, r.font.size, r.font.bold = FONT, Pt(14), True
        r.font.color.rgb = WHITE
    for i, row in enumerate(rows, start=1):
        strong = float(row[4].rstrip("σ")) >= 5.0
        for j, val in enumerate(row):
            c = tbl.cell(i, j)
            c.text = val
            r = c.text_frame.paragraphs[0].runs[0]
            r.font.name, r.font.size = FONT, Pt(14)
            r.font.color.rgb = INK
            if j >= 3 and strong:
                r.font.bold = True
                r.font.color.rgb = ORANGE
    add_text(s, Inches(7.7), Inches(1.8), Inches(5.1), Inches(4.8),
             bullets([
                 (0, "Same pixels, same painting: the frontal ratio "
                     "runs from 5.8 down to 0.63 across energy"),
                 (0, "Frontal repeatability ≤ 1.4 % - the tilt "
                     "moves the ratio up to 24σ above it"),
                 (0, "The shift decays monotonically with energy, "
                     "+9.5 % → +0.6 % - geometry, not drift"),
                 (0, "Ratios from the registered overlap region only: "
                     "a crop test showed full-frame ratios mix in a "
                     "field-of-view artifact"),
                 (0, "Four Pb lines share pixels and composition: a "
                     "pure energy dependence"),
             ]), space_after=12)
    add_notes(s, "3:00-4:15\n"
              "Table 1 in one look. Two messages: the channels are far "
              "from interchangeable (6x at calcium), and tilting "
              "produces a highly significant, monotonically decaying "
              "energy pattern. Mention the crop test: restricting to "
              "the registered overlap removed a field-of-view "
              "artifact - that is why these are the trustworthy "
              "numbers. The four lead lines are the internal "
              "control - same pixels, same composition, only energy "
              "differs.")

    # ---- 6 stage 1 ----------------------------------------------------
    s = blank_slide(prs)
    add_title(s, "Stage 1 - what separates the detectors")
    pic = s.shapes.add_picture(
        os.path.join(RESULTS, "geometry_fit.png"),
        Inches(0.55), Inches(1.75), width=Inches(6.4),
        height=Inches(5.12))
    pic.crop_right = 0.5
    add_text(s, Inches(7.35), Inches(1.9), Inches(5.4), Inches(4.8),
             bullets([
                 (0, "Three-parameter model: gain × differential "
                     "absorber × Si-thickness ratio"),
                 (0, "Reproduces R(E) over two orders of magnitude"),
                 (0, [("Fitted absorber: 973 ± 2 µm "
                       "Be-equivalent", 17, True, INK, False),
                      (" - 40-100× a real SDD window",
                       17, False, INK, False)]),
                 (0, [("⇒ ≈15-20 cm extra air path / "
                       "collimation in front of detector 19511",
                       17, True, ORANGE, False)]),
                 (0, "GAP"),
                 (0, [("An instrument diagnosis obtained from painting "
                       "scans alone.", 16, False, GRAY, True)]),
             ]), space_after=12)
    add_notes(s, "4:15-5:15\n"
              "The frontal curve is explained by detector properties "
              "alone. The surprise: the fitted low-energy absorber is "
              "the equivalent of a millimetre of beryllium - no window "
              "is that thick. It corresponds to 15-20 cm of air, so "
              "one detector must sit further away or behind a "
              "collimator. We learned that about the hardware without "
              "opening it.")

    # ---- 7 stage 2 ----------------------------------------------------
    s = blank_slide(prs)
    add_title(s, "Stage 2 - what the tilt reveals")
    pic = s.shapes.add_picture(
        os.path.join(RESULTS, "geometry_fit.png"),
        Inches(0.55), Inches(1.75), width=Inches(6.4),
        height=Inches(5.12))
    pic.crop_left = 0.5
    add_text(s, Inches(7.35), Inches(1.9), Inches(5.4), Inches(4.8),
             bullets([
                 (0, "Tilt shift of R: monotonic positive decay, "
                     "+9.5 % at 3.7 keV → +0.6 % at 14.8 keV"),
                 (0, "Thick-sample fluorescence model; the tilt moves "
                     "the effective take-off angles"),
                 (0, [("Lever arm s = 0.53 ± 0.10 ⇒ take-off "
                       "shift 4.1° ± 0.8° at this tilt",
                       17, True, INK, False)]),
                 (0, [("Gaussian-process check (no physics inside): "
                       "same shape, model within 1.1σ",
                       17, False, INK, False)]),
                 (0, [("⇒ the shape lives in the data, not in the "
                       "model", 17, True, ORANGE, False)]),
             ]), space_after=12)
    add_notes(s, "5:15-6:15\n"
              "The tilt-dependent part follows textbook fluorescence "
              "geometry. With a single tilt the individual take-off "
              "angles are not identifiable - what is, is how fast "
              "they move: about half a degree per degree of tilt. The "
              "grey band is a GP regression that knows no physics; "
              "the physical model stays inside it everywhere.")

    # ---- 8 self-measured angle ---------------------------------------
    s = blank_slide(prs)
    add_title(s, "The scan measures its own geometry")
    s.shapes.add_picture(
        os.path.join(RESULTS, "tilt_angle.png"),
        Inches(7.45), Inches(1.55), height=Inches(5.55))
    add_text(s, Inches(0.7), Inches(1.9), Inches(6.4), Inches(4.8),
             bullets([
                 (0, "A forward tilt compresses the image vertically "
                     "by cos θ"),
                 (0, "Register tilted ↔ frontal with independent "
                     "x/y scales: sy/sx = 1/cos θ - the unknown "
                     "step sizes cancel"),
                 (0, [("θ = 7.7° ± 1.0° ± 1.8° - an upper "
                       "bound (θ ≲ 8°)", 18, True, INK,
                       False)]),
                 (0, "A no-tilt control pair returns “5.3°”: "
                     "foreshortening at this angle sits near the "
                     "resolution floor"),
                 (0, "GAP"),
                 (0, [("Bonus: two “identical” frontal "
                       "sessions differ by 0.43 % in scale - a "
                       "measured positioning drift",
                       17, True, ORANGE, False)]),
             ]), space_after=12)
    add_notes(s, "6:15-7:15\n"
              "The mounting angle was never recorded - and it turns "
              "out we never needed it. The tilt compresses the image "
              "vertically, and the scan step sizes cancel in the "
              "scale ratio, so the angle falls out of registration: "
              "about 8 degrees. Be upfront: a control pair with no "
              "tilt returns five degrees, so we quote it as an upper "
              "bound until the instrument builder confirms. The "
              "precision floor itself is a result: it is the "
              "positioning drift between two nominally identical "
              "sessions.")

    # ---- 9 what this buys --------------------------------------------
    s = blank_slide(prs)
    add_title(s, "What the characterization buys",
              sub="fusion benchmark: cross-scan SNR on held-out pixels")
    add_text(s, Inches(0.7), Inches(1.45), SW - Inches(1.4), Inches(2.1),
             bullets([
                 (0, [("Flat-field map", 15, True, INK, False),
                      (": RMS ≈ 9 %, reproduced across scans "
                       "(r = 0.70) - and it is the scatter-artifact "
                       "geometry (r = 0.86)", 15, False, INK, False)]),
                 (0, [("Mounting-error price list", 15, True, INK,
                       False),
                      (": 0.50 %/° (Ca), 0.45 %/° (Ti), "
                       "≤ 0.13 %/° (Pb) - lower bounds",
                       15, False, INK, False)]),
                 (0, [("Fusion", 15, True, INK, False),
                      (": inverse-variance +0.9 % - the null is a "
                       "finding (Poisson-limited channels); ",
                       15, False, INK, False),
                      ("learned N2N +17.7 % mean SNR",
                       15, True, ORANGE, False),
                      (" (six lines from Fe up, Pb Ll +69 %; Ca/Ti "
                       "kept on the sum)", 15, False, INK, False)]),
             ]), space_after=8)
    s.shapes.add_picture(
        os.path.join(RESULTS, "fusion_benchmark.png"),
        Inches(1.17), Inches(3.45), width=Inches(11.0))
    add_notes(s, "7:15-7:45\n"
              "Everything here reuses the earlier numbers. The "
              "flat-field is reproducible and is literally the "
              "scatter-artifact geometry - one acquisition story. The "
              "weighting null is itself a finding: the channels are "
              "Poisson-limited, no scalar reweighting can win. The "
              "learned fusion goes beyond scalars: two detectors are "
              "two noise realizations of the same spectrum - "
              "Noise2Noise, no clean targets - and it gains 17.7 "
              "percent on pixels it never trained on. We quote it "
              "only where the absolute level is preserved; Ca and Ti "
              "stay on the sum. If asked: validation MSE does not "
              "rank these models like map SNR does - MSE rewards "
              "shrinking to the mean; hence the contrast guard.")

    # ---- 10 conclusions ----------------------------------------------
    s = blank_slide(prs)
    add_title(s, "Conclusions")
    add_text(s, Inches(0.9), Inches(1.7), SW - Inches(1.8), Inches(4.2),
             bullets([
                 (0, "Detector difference + one tilted scan = "
                     "geometry-resolved characterization, with no "
                     "dedicated calibration measurements"),
                 (0, "Per-element efficiency ratios (5.8 → 0.63); "
                     "detector vs geometry separated: take-off "
                     "response 0.5°/°, matrix scale Ec ≈ 3.6 keV"),
                 (0, [("Self-supervised fusion of the two channels: "
                       "+17.7 % SNR over summing on held-out pixels",
                       17, True, INK, False)]),
                 (0, [("The data self-report their geometry: tilt "
                       "≲ 8°, session drift 0.43 %, flat-field = "
                       "the artifact geometry", 17, True, ORANGE,
                       False)]),
                 (0, "GAP"),
                 (0, [("Data and code: github.com/JankoVukobratovic/"
                       "VincaInternship", 15, False, GRAY, False)]),
             ]), space_after=14)
    add_notes(s, "7:45-8:00\n"
              "One sentence to leave with: the signal everyone throws "
              "away is enough to characterize the instrument - "
              "including the angle nobody wrote down - and, once "
              "characterized, the two channels fuse into a better "
              "measurement than their sum. Thank you; happy to take "
              "questions.")

    out = os.path.join(OUTDIR, "dual_detector_talk.pptx")
    prs.save(out)
    print(f"Saved: {out}")
