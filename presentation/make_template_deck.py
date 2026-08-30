"""
make_template_deck.py
===============================================================================
Pour the dual-detector talk (presentation/dual_detector_talk.pptx, built by
make_slides.py) into the visual template of the ICETRAN 2026 deck
(poster/IcETRAN2026_final.pptx): same fonts (embedded Gothic A1 / Lovelo),
same greens, same slide furniture (title rule, outline boxes, three-column
motivation, five step cards, stat columns, two-column conclusions).

Nothing is drawn from scratch.  Every output slide is a clone of one of the
nine template slides -- shapes, background, picture fills, embedded fonts --
with the text replaced and the pictures swapped for the talk's figures.
Text keeps the run formatting of the box it lands in; a new box copies the
style of an existing one.  Pictures are placed with their own aspect ratio
(the template stretches a fill to the frame; here the frame is resized to
the picture).

Sources
    template   poster/IcETRAN2026_final.pptx
    figures    presentation/figs/geometry_schematic.png (make_slides.py)
               results/10264/prova1/element_maps.png
               results/detector_diff/geometry_fit.png   (split into 2 panels)
               results/detector_diff/tilt_angle.png
               results/detector_diff/fusion_benchmark.png
               results/detector_diff/fusion_showcase.png
    numbers    results/registration/overlap_ratios.csv
               results/detector_diff/geometry_fit.txt
    text       the talk slides and their speaker notes, transcribed below

Run from the project root:
    python presentation/make_template_deck.py
    -> presentation/dual_detector_talk_icetran_style.pptx
"""

import argparse
import copy
import csv
import os
import re
import sys

from lxml import etree
from PIL import Image
from pptx import Presentation
from pptx.oxml.ns import qn
from pptx.util import Cm, Pt

ROOT = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
os.chdir(ROOT)

TEMPLATE = os.path.join("poster", "IcETRAN2026_final.pptx")
OUT = os.path.join("presentation", "dual_detector_talk_icetran_style.pptx")
WORK = os.path.join("presentation", "figs", "template_deck")

RESULTS = os.path.join("results", "detector_diff")
GREEN = "456446"          # the template's text green
LIGHT = "9BBB9C"          # the template's title green (slide 1)
BLACK = "000000"

# --------------------------------------------------------------------------
# low-level helpers: clone, text, pictures, tables
# --------------------------------------------------------------------------
REL_ATTRS = (qn("r:embed"), qn("r:link"), qn("r:id"), qn("r:pict"))


def clone_slide(prs, src):
    """Append a copy of ``src`` (shapes, background, relationships)."""
    dst = prs.slides.add_slide(src.slide_layout)
    for sh in list(dst.shapes):
        sh._element.getparent().remove(sh._element)
    src_cSld, dst_cSld = src._element.cSld, dst._element.cSld
    bg = src_cSld.find(qn("p:bg"))
    if bg is not None:
        dst_cSld.insert(0, copy.deepcopy(bg))
    for el in src_cSld.spTree.iterchildren():
        if el.tag in (qn("p:nvGrpSpPr"), qn("p:grpSpPr")):
            continue
        dst_cSld.spTree.append(copy.deepcopy(el))
    rid_map = {}
    for el in dst_cSld.iter():
        for attr in REL_ATTRS:
            rid = el.get(attr)
            if not rid:
                continue
            if rid not in rid_map:
                rel = src.part.rels[rid]
                if rel.is_external:
                    rid_map[rid] = dst.part.rels.get_or_add_ext_rel(
                        rel.reltype, rel.target_ref)
                else:
                    rid_map[rid] = dst.part.relate_to(rel.target_part,
                                                      rel.reltype)
            el.set(attr, rid_map[rid])
    return dst


def delete_slide(prs, slide):
    sldIdLst = prs.slides._sldIdLst
    for sldId in list(sldIdLst):
        if prs.slides.part.related_part(sldId.rId) is slide.part:
            prs.part.drop_rel(sldId.rId)
            sldIdLst.remove(sldId)
            return
    raise KeyError("slide not found")


def shape(slide, name):
    for sh in slide.shapes:
        if sh.name == name:
            return sh
    raise KeyError("%s has no shape %r" % (slide, name))


def remove(slide, *names):
    for n in names:
        el = shape(slide, n)._element
        el.getparent().remove(el)


def _style_of(txBody):
    """(pPr, rPr) of the first run in a txBody, as deep copies."""
    p0 = txBody.find(qn("a:p"))
    pPr = p0.find(qn("a:pPr"))
    r0 = p0.find(qn("a:r"))
    rPr = r0.find(qn("a:rPr")) if r0 is not None else None
    if rPr is None:
        rPr = p0.find(qn("a:endParaRPr"))
    return (copy.deepcopy(pPr) if pPr is not None else None,
            copy.deepcopy(rPr) if rPr is not None else etree.Element(qn("a:rPr")))


def _set_color(rPr, hexcolor):
    for old in rPr.findall(qn("a:solidFill")):
        rPr.remove(old)
    fill = etree.Element(qn("a:solidFill"))
    clr = etree.SubElement(fill, qn("a:srgbClr"))
    clr.set("val", hexcolor)
    rPr.insert(0, fill)


def _strip_bullet(pPr):
    for tag in ("a:buFont", "a:buChar", "a:buNone", "a:buAutoNum"):
        for el in pPr.findall(qn(tag)):
            pPr.remove(el)
    etree.SubElement(pPr, qn("a:buNone"))
    pPr.set("marL", "0")
    pPr.set("indent", "0")


def _set_bullet(pPr, char, level):
    for tag in ("a:buFont", "a:buChar", "a:buNone", "a:buAutoNum"):
        for el in pPr.findall(qn(tag)):
            pPr.remove(el)
    marL = 360000 + 360000 * level
    pPr.set("marL", str(marL))
    pPr.set("indent", str(-300000))
    bf = etree.SubElement(pPr, qn("a:buFont"))
    bf.set("typeface", "Arial")
    bc = etree.SubElement(pPr, qn("a:buChar"))
    bc.set("char", char)


def set_paras(target, items, style=None):
    """Replace the paragraphs of ``target`` (a shape or a table cell).

    ``items``: strings or (string, options) pairs.  Options: size (pt),
    bold, color (hex), bullet (True / False / a bullet character), level,
    space_after (pt), align ('l', 'ctr', 'r').  The paragraph and run
    formatting of the first existing paragraph is the base, or ``style``
    = (pPr, rPr) taken from another shape via style_of().
    """
    txBody = target.text_frame._txBody
    pPr_t, rPr_t = style if style is not None else _style_of(txBody)
    for p in txBody.findall(qn("a:p")):
        txBody.remove(p)
    for it in items:
        text, opt = (it, {}) if isinstance(it, str) else it
        p = etree.SubElement(txBody, qn("a:p"))
        pPr = copy.deepcopy(pPr_t) if pPr_t is not None else etree.Element(qn("a:pPr"))
        p.append(pPr)
        if "bullet" in opt:
            if opt["bullet"] is False:
                _strip_bullet(pPr)
            else:
                _set_bullet(pPr, "•" if opt["bullet"] is True else opt["bullet"],
                            opt.get("level", 0))
        if "align" in opt:
            pPr.set("algn", opt["align"])
        if "space_after" in opt:
            for el in pPr.findall(qn("a:spcAft")):
                pPr.remove(el)
            sa = etree.SubElement(pPr, qn("a:spcAft"))
            sp = etree.SubElement(sa, qn("a:spcPts"))
            sp.set("val", str(int(opt["space_after"] * 100)))
        if "line" in opt:
            for el in pPr.findall(qn("a:lnSpc")):
                pPr.remove(el)
            ls = etree.Element(qn("a:lnSpc"))
            sp = etree.SubElement(ls, qn("a:spcPts"))
            sp.set("val", str(int(opt["line"] * 100)))
            pPr.insert(0, ls)
        r = etree.SubElement(p, qn("a:r"))
        rPr = copy.deepcopy(rPr_t)
        rPr.tag = qn("a:rPr")
        r.append(rPr)
        if "size" in opt:
            rPr.set("sz", str(int(opt["size"] * 100)))
        if "bold" in opt:
            rPr.set("b", "true" if opt["bold"] else "false")
        if "color" in opt:
            _set_color(rPr, opt["color"])
        t = etree.SubElement(r, qn("a:t"))
        t.text = text


def style_of(sh):
    return _style_of(sh.text_frame._txBody)


def add_text(slide, x, y, w, h, items, style):
    tb = slide.shapes.add_textbox(Cm(x), Cm(y), Cm(w), Cm(h))
    bodyPr = tb.text_frame._txBody.find(qn("a:bodyPr"))
    for k in ("lIns", "tIns", "rIns", "bIns"):
        bodyPr.set(k, "0")
    bodyPr.set("wrap", "square")
    set_paras(tb, items, style)
    return tb


def set_picture(slide, sh, path, box=None, fit="inside", anchor="lt"):
    """Swap the picture fill of a freeform/picture for ``path``.

    The frame is resized to the picture's aspect ratio inside ``box`` =
    (x, y, w, h) in cm (default: the shape's own box).  ``anchor`` picks
    which edges of the box the picture keeps: 'l'/'r'/'c' and 't'/'b'/'m'.
    """
    _, rId = slide.part.get_or_add_image_part(path)
    blip = sh._element.find(".//" + qn("a:blip"))
    blip.set(qn("r:embed"), rId)
    for child in list(blip):
        blip.remove(child)          # alphaModFix etc. belong to the old fill
    fr = sh._element.find(".//" + qn("a:fillRect"))
    if fr is not None:
        for k in ("l", "t", "r", "b"):
            fr.set(k, "0")
    if box is None:
        box = (sh.left / 360000, sh.top / 360000,
               sh.width / 360000, sh.height / 360000)
    x, y, w, h = box
    with Image.open(path) as im:
        aspect = im.height / im.width
    if fit == "width" or (fit == "inside" and w * aspect <= h):
        nw, nh = w, w * aspect
    else:
        nh, nw = h, h / aspect
    nx = x if anchor[0] == "l" else (x + w - nw if anchor[0] == "r" else x + (w - nw) / 2)
    ny = y if anchor[1] == "t" else (y + h - nh if anchor[1] == "b" else y + (h - nh) / 2)
    sh.left, sh.top, sh.width, sh.height = Cm(nx), Cm(ny), Cm(nw), Cm(nh)
    return sh


def resize_table(sh, n_rows, n_cols, row_h_cm, col_w_cm=None):
    tbl = sh.table._tbl
    grid = tbl.tblGrid
    cols = grid.findall(qn("a:gridCol"))
    while len(cols) < n_cols:
        grid.append(copy.deepcopy(cols[-1]))
        cols = grid.findall(qn("a:gridCol"))
    while len(cols) > n_cols:
        grid.remove(cols[-1])
        cols = grid.findall(qn("a:gridCol"))
    rows = tbl.findall(qn("a:tr"))
    while len(rows) < n_rows:
        tbl.append(copy.deepcopy(rows[-1]))
        rows = tbl.findall(qn("a:tr"))
    while len(rows) > n_rows:
        tbl.remove(rows[-1])
        rows = tbl.findall(qn("a:tr"))
    for tr in rows:
        tcs = tr.findall(qn("a:tc"))
        while len(tcs) < n_cols:
            tr.append(copy.deepcopy(tcs[-1]))
            tcs = tr.findall(qn("a:tc"))
        while len(tcs) > n_cols:
            tr.remove(tcs[-1])
            tcs = tr.findall(qn("a:tc"))
        tr.set("h", str(Cm(row_h_cm)))
    widths = col_w_cm or [sh.width / 360000 / n_cols] * n_cols
    for gc, wcm in zip(grid.findall(qn("a:gridCol")), widths):
        gc.set("w", str(Cm(wcm)))
    sh.width = Cm(sum(widths))
    sh.height = Cm(row_h_cm * n_rows)


def fill_table(sh, rows, size=None, bold_first_row=False, align=None):
    for r, row in enumerate(rows):
        for c, text in enumerate(row):
            cell = sh.table.cell(r, c)
            opt = {}
            if size:
                opt["size"] = size
            if bold_first_row and r == 0:
                opt["bold"] = True
            if align:
                opt["align"] = align[c]
            set_paras(cell, [(text, opt)])
            cell.margin_top = cell.margin_bottom = Cm(0.1)


def set_notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text


def title(slide, name, text, size=54, width=45.0):
    sh = shape(slide, name)
    sh.width = Cm(width)
    set_paras(sh, [(text, {"size": size})])


# --------------------------------------------------------------------------
# figure preparation
# --------------------------------------------------------------------------
def trim_white(im, pad=8):
    from PIL import ImageChops
    bg = Image.new(im.mode, im.size, (255, 255, 255))
    diff = ImageChops.difference(im.convert("RGB"), bg.convert("RGB"))
    bbox = diff.getbbox()
    if bbox is None:
        return im
    x0, y0, x1, y1 = bbox
    return im.crop((max(0, x0 - pad), max(0, y0 - pad),
                    min(im.width, x1 + pad), min(im.height, y1 + pad)))


def split_panels(src, left_out, right_out):
    im = Image.open(src).convert("RGB")
    half = im.width // 2
    trim_white(im.crop((0, 0, half, im.height))).save(left_out)
    trim_white(im.crop((half, 0, im.width, im.height))).save(right_out)


def render_formula(tex, path, size=30):
    import matplotlib
    matplotlib.use("Agg")
    import matplotlib.pyplot as plt
    with plt.rc_context({"mathtext.fontset": "cm", "font.size": size}):
        fig = plt.figure(figsize=(12, 2.4))
        fig.text(0.5, 0.5, tex, ha="center", va="center", color="#" + GREEN)
        fig.savefig(path, dpi=300, transparent=True, bbox_inches="tight",
                    pad_inches=0.08)
        plt.close(fig)


def prepare_figures():
    os.makedirs(WORK, exist_ok=True)
    figs = {
        "schematic": os.path.join("presentation", "figs", "geometry_schematic.png"),
        "maps": os.path.join("results", "10264", "prova1", "element_maps.png"),
        "tilt": os.path.join(RESULTS, "tilt_angle.png"),
        "benchmark": os.path.join(RESULTS, "fusion_benchmark.png"),
        "showcase": os.path.join(RESULTS, "fusion_showcase.png"),
    }
    for k, f in figs.items():
        if not os.path.exists(f):
            sys.exit("ERROR: missing figure %s (%s)" % (k, f))
    figs["stage1"] = os.path.join(WORK, "geometry_fit_stage1.png")
    figs["stage2"] = os.path.join(WORK, "geometry_fit_stage2.png")
    split_panels(os.path.join(RESULTS, "geometry_fit.png"), figs["stage1"],
                 figs["stage2"])
    figs["eq1"] = os.path.join(WORK, "eq_stage1.png")
    figs["eq2"] = os.path.join(WORK, "eq_stage2.png")
    render_formula(
        r"$R_{\mathrm{det}}(E) \;=\; k \; e^{\,\rho_{\mathrm{Be}}\mu_{\mathrm{Be}}(E)\,d}\;"
        r"\dfrac{1 - e^{-\rho_{\mathrm{Si}}\mu_{\mathrm{Si}}(E)\,t_1}}"
        r"{1 - e^{-\rho_{\mathrm{Si}}\mu_{\mathrm{Si}}(E)\,t_2}}$",
        figs["eq1"])
    render_formula(
        r"$\delta(E) = (1+c)\,\dfrac{g(\theta)}{g(0)} - 1,\qquad "
        r"g(\theta) = \dfrac{1/\sin\varphi + x/\sin\psi_-}{1/\sin\varphi + x/\sin\psi_+},"
        r"\qquad x = (E/E_c)^{-3}$",
        figs["eq2"], size=26)
    return figs


def ratio_rows():
    """(line, R frontal, tilt shift) for the reliable lines, from script 08."""
    label = {"Ca": "Ca Kα", "Ti": "Ti Kα", "Fe": "Fe Kα",
             "Cu": "Cu Kα", "PbLl": "Pb Lℓ", "PbLa": "Pb Lα",
             "PbLb": "Pb Lβ", "PbLg": "Pb Lγ"}
    rows = []
    with open(os.path.join("results", "registration", "overlap_ratios.csv"),
              newline="") as f:
        for r in csv.DictReader(f):
            if r["reliable"] != "True":
                continue
            rows.append((float(r["kev"]), label.get(r["element"], r["element"]),
                         float(r["R_frontal_overlap"]),
                         float(r["tilt_overlap_pct"])))
    rows.sort()
    return [(name, "%.2f" % R, "%+.1f %%" % d) for _, name, R, d in rows]


# --------------------------------------------------------------------------
# the deck
# --------------------------------------------------------------------------
def build(template, out):
    figs = prepare_figures()
    prs = Presentation(template)
    T = list(prs.slides)              # the nine template slides
    if len(T) != 9:
        sys.exit("ERROR: expected 9 template slides, found %d" % len(T))
    t_title, t_outline, t_columns, t_picture, t_cards, t_formula, \
        t_results, t_stats, t_conclusions = T

    # styles borrowed from the template
    body_style = style_of(shape(t_columns, "TextBox 4"))        # 30 pt bullets
    head_style = style_of(shape(t_columns, "TextBox 3"))        # bold column head
    small_style = style_of(shape(t_results, "TextBox 29"))      # 21.8 pt note

    def bullets(items, size=24, line=None):
        out = []
        for it in items:
            text, opt = (it, {}) if isinstance(it, str) else it
            o = {"size": size, "bullet": "–" if opt.get("level") else True,
                 "level": opt.get("level", 0), "space_after": 6}
            if line:
                o["line"] = line
            o.update({k: v for k, v in opt.items() if k in ("bold", "color", "size")})
            out.append((text, o))
        return out

    # ---- 1  title --------------------------------------------------------
    s = clone_slide(prs, t_title)
    set_paras(shape(s, "TextBox 6"), [
        "Geometry-resolved Characterization",
        "of a Dual-detector MA-XRF Scanner"])
    set_paras(shape(s, "TextBox 7"), [
        "Dimitrije Pešić · Janko Vukobratović · "
        "Aleksandra Stojanović · Giulia Ristori · Stefano Ridolfi "
        "· Maja Gajić-Kvaščev · Goran Kvaščev",
        ("University of Belgrade, School of Electrical Engineering · "
         "Ars Mensurae, Rome · IDArtScience, Rome · "
         "Vinča Institute of Nuclear Sciences / VINARH", {"size": 16})])
    set_notes(s, "0:00-0:15 Good morning. I will show how a dual-detector MA-XRF "
                 "scanner can be characterized - efficiencies, geometry, even the "
                 "mounting angle - using nothing but routine scans of a painting.")

    # ---- 2  outline ------------------------------------------------------
    s = clone_slide(prs, t_outline)
    set_paras(shape(s, "TextBox 3"), ["Outline"])
    for name, text in (("TextBox 4", "Motivation"), ("TextBox 5", "Method"),
                       ("TextBox 6", "Data"), ("TextBox 8", "Results"),
                       ("TextBox 7", "Geometry & fusion"),
                       ("TextBox 9", "Conclusions")):
        set_paras(shape(s, name), [text])

    # ---- 3  motivation, three columns -----------------------------------
    s = clone_slide(prs, t_columns)
    title(s, "TextBox 2", "Two detectors, one discarded signal")
    set_paras(shape(s, "TextBox 3"), ["Problem"])
    set_paras(shape(s, "TextBox 8"), ["Insight"])
    set_paras(shape(s, "TextBox 9"), ["Our approach"])
    set_paras(shape(s, "TextBox 4"), bullets([
        "MA-XRF scans a painting pixel by pixel; XRF lines give element maps "
        "(Ca, Fe, Pb, ...)",
        "Many scanners carry two SDD detectors; their signals are summed for "
        "signal-to-noise",
        "The difference between the two channels is treated as noise and "
        "thrown away"], size=24))
    set_paras(shape(s, "TextBox 5"), bullets([
        ("That discarded difference is a calibration signal.", {"bold": True}),
        "The two heads see the same spot from opposite take-off angles, so "
        "their ratio carries the detector response and the acquisition "
        "geometry",
        "Everything that follows is extracted from that difference"], size=24))
    set_paras(shape(s, "TextBox 6"), bullets([
        "One extra tilted scan separates the detector part from the "
        "geometric part",
        "Two-stage physical model, no standards, no extra hardware, no "
        "dedicated calibration measurements",
        "Once characterized, the two channels fuse into a better measurement "
        "than their sum (self-supervised Noise2Noise)"], size=24))
    set_notes(s, "0:15-1:15 Set the scene: MA-XRF, two detectors, summing is "
                 "standard practice. The premise of this work: the part everyone "
                 "throws away carries quantitative information about the "
                 "instrument itself. Everything that follows is extracted from "
                 "that difference.")

    # ---- 4  method: the tilt splits the physics --------------------------
    s = clone_slide(prs, t_picture)
    title(s, "TextBox 6", "One extra tilted scan splits the physics")
    remove(s, "Freeform 4", "Freeform 5")
    set_picture(s, shape(s, "Freeform 2"), figs["schematic"],
                box=(27.5, 7.6, 20.4, 18.5), anchor="rt")
    add_text(s, 2.9, 8.0, 23.0, 18.0, bullets([
        "Scan the same painting twice: frontal, and with the canvas tilted "
        "forward",
        "Tilting changes the photon paths to each detector; the detectors "
        "stay the same",
        "The per-element ratio R(E) = det1 / det2 then separates into:",
        ("detector part: tilt-invariant (absorbers, Si thickness, gain)",
         {"level": 1, "bold": True}),
        ("geometric part: tilt-dependent, energy-structured",
         {"level": 1, "bold": True}),
        ("“with no dedicated calibration measurements”",
         {"bold": True})], size=24), body_style)
    set_notes(s, "1:15-2:15 The one-slide method. The tilt is the control knob: "
                 "it moves only the geometry. Comparing tilted vs frontal ratios "
                 "cleanly separates detector properties from acquisition "
                 "geometry. Quote the abstract: no dedicated calibration "
                 "measurements - no reference targets, no monochromatic sources.")

    # ---- 5  instrument and data -----------------------------------------
    s = clone_slide(prs, t_picture)
    title(s, "TextBox 6", "Instrument and data")
    remove(s, "Freeform 4", "Freeform 5")
    add_text(s, 2.9, 7.6, 45.0, 6.5, bullets([
        "Canvas copy of the Creation of Adam; Ars Mensurae MA-XRF scanner, "
        "two SDDs (s/n 10264, 19511), 40 kV / 50 µA",
        "Three scans: two frontal, 7 days apart (repeatability baseline), "
        "one tilted forward",
        "8 usable lines from Ca Kα (3.7 keV) to Pb Lγ (14.8 keV); "
        "bootstrap uncertainties"], size=24), body_style)
    set_picture(s, shape(s, "Freeform 2"), figs["maps"],
                box=(4.4, 14.4, 42.0, 13.4), fit="width", anchor="ct")
    set_notes(s, "2:15-3:00 The test object is a canvas copy of the Creation of "
                 "Adam - the element maps show the two hands. Three routine "
                 "scans, nothing else. The two frontal scans give the "
                 "repeatability floor; the tilted one is the geometry probe.")

    # ---- 6  pipeline, five cards ----------------------------------------
    s = clone_slide(prs, t_cards)
    title(s, "TextBox 3", "The pipeline in five steps")
    cards = [
        ("TextBox 12", "TextBox 13", "Registration",
         "tilted onto frontal; affine, NCC 0.965; overlap only"),
        ("TextBox 22", "TextBox 23", "Line ratios R(E)",
         "8 lines, Ca Kα to Pb Lγ; bootstrap σ"),
        ("TextBox 32", "TextBox 33", "Stage 1: detector",
         "gain × absorber × Si thickness; frontal curve"),
        ("TextBox 42", "TextBox 43", "Stage 2: geometry",
         "thick-sample model; tilt shift of R"),
        ("TextBox 52", "TextBox 53", "Fusion",
         "Noise2Noise 1D U-Net; R(E)-weighted loss"),
    ]
    for head, body, h, b in cards:
        sh = shape(s, head)
        sh.width = Cm(6.6)
        set_paras(sh, [(h, {"align": "ctr"})])
        sh = shape(s, body)
        sh.width = Cm(6.6)
        sh.left = shape(s, head).left
        set_paras(sh, [(b, {"align": "ctr", "size": 20})])
    set_notes(s, "The whole chain, left to right: register the tilted scan onto "
                 "the frontal one, form the per-line ratios on the overlap, fit "
                 "the frontal curve with the detector model, fit the tilt shift "
                 "with the geometry model, and finally use the characterized "
                 "variance structure to fuse the two channels.")

    # ---- 7  results: the two channels disagree --------------------------
    s = clone_slide(prs, t_results)
    title(s, "TextBox 27", "The two channels disagree, up to sixfold")
    sh = shape(s, "TextBox 28")
    sh.width = Cm(24.0)
    set_paras(sh, ["Frontal ratio and tilt shift per line"])
    sh = shape(s, "TextBox 29")
    sh.width = Cm(24.0)
    set_paras(sh, ["registered overlap, bootstrap σ; R = det 10264 / det 19511"])
    tbl = shape(s, "Table 3")
    rows = [("line", "R frontal", "tilt shift")] + ratio_rows()
    resize_table(tbl, len(rows), 3, 1.45, [5.4, 5.3, 5.3])
    tbl.top = Cm(9.9)
    fill_table(tbl, rows, size=16, bold_first_row=True, align=("l", "ctr", "ctr"))
    y_after = 9.9 + 1.45 * len(rows) + 0.5
    sh = shape(s, "TextBox 4")
    sh.top = Cm(y_after)
    set_paras(sh, ["REGISTERED OVERLAP ONLY"])
    t5 = shape(s, "Table 5")
    t5.top = Cm(y_after + 1.5)
    fill_table(t5, [("affine", "NCC 0.965", "crop-tested")], size=14,
               align=("ctr", "ctr", "ctr"))
    # key numbers: label and value merged into one box each
    set_paras(shape(s, "TextBox 31"), ["KEY NUMBERS"])
    for lab, val, text in (("TextBox 18", "TextBox 22", "Ca Kα:  R = 5.8"),
                           ("TextBox 19", "TextBox 23", "Pb Lγ:  R = 0.63"),
                           ("TextBox 20", "TextBox 24", "repeatability ≤ 1.4 %"),
                           ("TextBox 21", "TextBox 25", "tilt shift up to 24σ")):
        sh = shape(s, lab)
        sh.width = Cm(12.0)
        set_paras(sh, [(text, {"size": 24})])
        remove(s, val)
    remove(s, "Freeform 26", "Freeform 30")
    sh = shape(s, "TextBox 32")
    sh.top = Cm(14.6)
    sh.width = Cm(26.5)
    set_paras(sh, [
        ("SHIFT DECAYS WITH ENERGY", {"size": 24, "bold": True, "space_after": 6}),
        ("+9.5 % at 3.7 keV → +0.6 % at 14.8 keV, monotonic: geometry, "
         "not drift", {"size": 22, "bold": False, "space_after": 6}),
        ("Four Pb lines share pixels and composition: a pure energy dependence",
         {"size": 22, "bold": False, "space_after": 6}),
        ("A crop test showed full-frame ratios mix in a field-of-view "
         "artifact; the overlap numbers are the trustworthy ones",
         {"size": 22, "bold": False})])
    set_notes(s, "3:00-4:15 Table 1 in one look. Two messages: the channels are "
                 "far from interchangeable (6x at calcium), and tilting produces "
                 "a highly significant, monotonically decaying energy pattern. "
                 "Mention the crop test: restricting to the registered overlap "
                 "removed a field-of-view artifact - that is why these are the "
                 "trustworthy numbers. The four lead lines are the internal "
                 "control - same pixels, same composition, only energy differs.")

    # ---- 8 / 9  stage 1 and stage 2 -------------------------------------
    def stage_slide(ttl, head, sub, columns, panel, eq, closing, notes):
        s = clone_slide(prs, t_stats)
        title(s, "TextBox 5", ttl)
        set_paras(shape(s, "TextBox 6"), [head])
        sh = shape(s, "TextBox 7")
        sh.width = Cm(25.0)
        set_paras(sh, [(sub, {"size": 22})])
        for name, lines, x in zip(("TextBox 8", "TextBox 9", "TextBox 10"),
                                  columns, (2.9, 11.5, 20.1)):
            sh = shape(s, name)
            sh.left = Cm(x)
            sh.top = Cm(13.2)
            sh.width = Cm(8.4)
            items = [(lines[0], {"size": 24, "bold": True, "space_after": 4}),
                     (lines[1], {"size": 28, "bold": True, "space_after": 4})]
            items += [(ln, {"size": 18, "bold": False, "space_after": 2})
                      for ln in lines[2:]]
            set_paras(sh, items)
        remove(s, "TextBox 11", "TextBox 12")
        set_picture(s, shape(s, "Freeform 4"), panel,
                    box=(28.6, 7.2, 19.3, 13.4), anchor="rt")
        set_picture(s, shape(s, "Freeform 3"), eq,
                    box=(28.6, 21.2, 19.3, 5.6), anchor="rm")
        add_text(s, 2.9, 23.8, 25.0, 3.5,
                 [(closing, {"size": 22, "bold": True, "bullet": False})],
                 body_style)
        set_notes(s, notes)
        return s

    stage_slide(
        "Stage 1: what separates the detectors", "DETECTOR MODEL",
        "Three-parameter model, gain × differential absorber × "
        "Si-thickness ratio, reproduces R(E) over two orders of magnitude",
        [("Absorber", "973 ± 2 µm", "Be-equivalent",
          "40-100× a real SDD window"),
         ("Air path", "≈ 15-20 cm", "extra path / collimation",
          "in front of detector 19511"),
         ("Silicon", "302 ± 1 µm", "active thickness of 10264",
          "(19511 fixed at 500 µm)")],
        figs["stage1"], figs["eq1"],
        "An instrument diagnosis obtained from painting scans alone.",
        "4:15-5:15 The frontal curve is explained by detector properties "
        "alone. The surprise: the fitted low-energy absorber is the equivalent "
        "of a millimetre of beryllium - no window is that thick. It corresponds "
        "to 15-20 cm of air, so one detector must sit further away or behind a "
        "collimator. We learned that about the hardware without opening it. "
        "Manufacturer specs since received (X-123SDD family): nominal window "
        "8-12.5 um Be and 500 um Si - the fitted reference matches the factory "
        "value; snout extenders up to 9 in exist, so a long-extender mounting "
        "of 19511 is the concrete candidate.")

    stage_slide(
        "Stage 2: what the tilt reveals", "GEOMETRIC MODEL",
        "Thick-sample fluorescence; the tilt moves the effective take-off "
        "angles antisymmetrically, ψ± = ψ₀ ± sθ",
        [("Lever arm", "s = 0.53 ± 0.10", "deg per deg of tilt",
          "take-off shift 4.1° ± 0.8°"),
         ("Matrix scale", "Ec = 3.6 ± 0.5 keV", "x = (E/Ec)⁻³",
          "offset c = +0.012 ± 0.002"),
         ("GP check", "within 1.1σ", "no physics inside",
          "same shape as the data")],
        figs["stage2"], figs["eq2"],
        "The shape lives in the data, not in the model: +9.5 % at 3.7 keV "
        "→ +0.6 % at 14.8 keV.",
        "5:15-6:15 The tilt-dependent part follows textbook fluorescence "
        "geometry. With a single tilt the individual take-off angles are not "
        "identifiable - what is, is how fast they move: about half a degree "
        "per degree of tilt. The grey band is a GP regression that knows no "
        "physics; the physical model stays inside it everywhere.")

    # ---- 10  the scan measures its own geometry -------------------------
    s = clone_slide(prs, t_picture)
    title(s, "TextBox 6", "The scan measures its own geometry")
    remove(s, "Freeform 4", "Freeform 5")
    set_picture(s, shape(s, "Freeform 2"), figs["tilt"],
                box=(29.0, 7.4, 18.9, 18.9), anchor="rt")
    add_text(s, 2.9, 7.8, 25.0, 19.0, bullets([
        "A forward tilt compresses the image vertically by cos θ",
        "Register tilted ↔ frontal with independent x/y scales: "
        "sy/sx = 1/cos θ; the unknown step sizes cancel",
        ("θ = 7.7° ± 1.0° ± 1.8°, an upper bound "
         "(θ ≲ 8°)", {"bold": True}),
        "A no-tilt control pair returns “5.3°”: foreshortening "
        "at this angle sits near the resolution floor",
        ("Bonus: two “identical” frontal sessions differ by 0.43 % "
         "in scale, a measured positioning drift", {"bold": True})],
        size=22), body_style)
    set_notes(s, "6:15-7:15 The mounting angle was never recorded - and it turns "
                 "out we never needed it. The tilt compresses the image "
                 "vertically, and the scan step sizes cancel in the scale ratio, "
                 "so the angle falls out of registration: about 8 degrees. Be "
                 "upfront: a control pair with no tilt returns five degrees, so "
                 "we quote it as an upper bound until the instrument builder "
                 "confirms. The precision floor itself is a result: it is the "
                 "positioning drift between two nominally identical sessions.")

    # ---- 11  what the characterization buys -----------------------------
    s = clone_slide(prs, t_formula)
    title(s, "TextBox 3", "What the characterization buys")
    remove(s, "Group 4")
    set_picture(s, shape(s, "Picture 18"), figs["benchmark"],
                box=(2.9, 7.6, 25.5, 9.5), fit="width", anchor="lt")
    sh = shape(s, "TextBox 20")
    sh.width = Cm(19.2)
    set_paras(sh, ["ERROR BUDGET AND GAINS"])
    tbl = shape(s, "Table 19")
    rows = [("quantity", "value", "note"),
            ("flat-field RMS", "≈ 9 %", "r = 0.70 across scans"),
            ("scatter-artifact match", "r = 0.86", "one acquisition story"),
            ("Ca per degree", "0.50 %/°", "lower bound"),
            ("Ti per degree", "0.45 %/°", "lower bound"),
            ("Pb per degree", "≤ 0.13 %/°", "lower bound"),
            ("inverse-variance fusion", "+0.9 %", "Poisson-limited null"),
            ("learned N2N fusion", "+17.7 %", "held-out mean SNR")]
    resize_table(tbl, len(rows), 3, 1.55, [7.4, 4.6, 7.2])
    tbl.top = Cm(10.6)
    fill_table(tbl, rows, size=15, bold_first_row=True, align=("l", "ctr", "l"))
    add_text(s, 2.9, 17.8, 25.5, 10.0, bullets([
        "Fusion benchmark: cross-scan SNR on held-out pixels. The weighting "
        "null is a finding: Poisson-limited channels, no scalar reweighting "
        "can win",
        "Learned N2N: +17.7 % mean SNR, six lines from Fe up, Pb Lℓ +69 %; "
        "Ca/Ti kept on the sum",
        "Canvas relief from one scan: local-slope map from the per-pixel ratio "
        "residuals, r = 0.73 across scans, tilt-law consistent "
        "(χ²/dof 0.85)"], size=20), body_style)
    set_notes(s, "7:15-7:45 Everything here reuses the earlier numbers. The "
                 "flat-field is reproducible and is literally the "
                 "scatter-artifact geometry - one acquisition story. The "
                 "weighting null is itself a finding: the channels are "
                 "Poisson-limited, no scalar reweighting can win. The learned "
                 "fusion goes beyond scalars: two detectors are two noise "
                 "realizations of the same spectrum - Noise2Noise, no clean "
                 "targets - and it gains 17.7 percent on pixels it never trained "
                 "on. We quote it only where the absolute level is preserved; Ca "
                 "and Ti stay on the sum.")

    # ---- 12  conclusions -------------------------------------------------
    s = clone_slide(prs, t_conclusions)
    title(s, "TextBox 3", "Conclusions")
    set_paras(shape(s, "TextBox 4"), ["WHAT WE SHOWED"])
    set_paras(shape(s, "TextBox 5"), bullets([
        "Detector difference + one tilted scan = geometry-resolved "
        "characterization, with no dedicated calibration measurements",
        "Per-element efficiency ratios (5.8 → 0.63); detector vs geometry "
        "separated: take-off response 0.5°/°, matrix scale "
        "Ec ≈ 3.6 keV",
        ("Self-supervised fusion of the two channels: +17.7 % SNR over "
         "summing on held-out pixels", {"bold": True}),
        ("The data self-report their geometry: tilt ≲ 8°, session "
         "drift 0.43 %, flat-field = the artifact geometry", {"bold": True})],
        size=22))
    sh = shape(s, "TextBox 6")
    sh.width = Cm(14.0)
    set_paras(sh, ["TAKE-HOME"])
    set_paras(shape(s, "TextBox 7"), bullets([
        "The signal everyone throws away is enough to characterize the "
        "instrument, including the angle nobody wrote down",
        "Once characterized, the two channels fuse into a better measurement "
        "than their sum",
        "Field MA-XRF scanners can be characterized from their own routine "
        "scans, with mounting geometry in the error budget",
        "Data and code: github.com/JankoVukobratovic/VincaInternship"],
        size=22))
    set_notes(s, "7:45-8:00 One sentence to leave with: the signal everyone "
                 "throws away is enough to characterize the instrument - "
                 "including the angle nobody wrote down - and, once "
                 "characterized, the two channels fuse into a better measurement "
                 "than their sum. Thank you; happy to take questions.")

    # ---- 13  backup: fusion on the maps ---------------------------------
    s = clone_slide(prs, t_picture)
    title(s, "TextBox 6", "Backup: fusion seen on the maps")
    remove(s, "Freeform 4", "Freeform 5")
    set_picture(s, shape(s, "Freeform 2"), figs["showcase"],
                box=(19.5, 7.4, 28.4, 19.5), anchor="rt")
    add_text(s, 2.9, 7.8, 15.5, 19.0, bullets([
        "Cross-scan noise panels: summed vs learned, held-out SNR annotated",
        "Pb Lℓ 5.4 → 9.1, Fe 10.6 → 13.7 on pixels the network "
        "never trained on",
        "Same signal structure, visibly calmer noise",
        "Contrast guard cv = 0.91-0.99: not blurring"], size=20), body_style)
    set_notes(s, "BACKUP - not part of the 8:00. If asked how the fusion gain "
                 "looks in practice: same signal structure, visibly calmer noise "
                 "panels; Pb Ll 5.4 to 9.1, Fe 10.6 to 13.7 on pixels the "
                 "network never trained on. Contrast guard cv = 0.91-0.99: not "
                 "blurring.")

    # ---- 14  backup: why not the individual angles ----------------------
    s = clone_slide(prs, t_conclusions)
    title(s, "TextBox 3", "Backup: why not the individual take-off angles?")
    sh = shape(s, "TextBox 4")
    sh.width = Cm(21.0)
    set_paras(sh, ["NOT IDENTIFIABLE WITH ONE TILT"])
    set_paras(shape(s, "TextBox 5"), bullets([
        "Fitting both take-off angles directly returns ~1e5-degree "
        "uncertainties: with a single tilt the shift is first-order in the "
        "response, second-order in the asymmetry",
        "Same story for the tilt angle itself: foreshortening at 7.7° "
        "sits near the registration floor (a no-tilt control pair returns "
        "“5.3°”), hence the upper bound"], size=22))
    sh = shape(s, "TextBox 6")
    sh.width = Cm(21.0)
    set_paras(sh, ["WHAT IS IDENTIFIABLE"])
    set_paras(shape(s, "TextBox 7"), bullets([
        ("The lever arm s = 0.53 ± 0.10 deg/deg: how fast the effective "
         "angles move per degree of tilt", {"bold": True}),
        "With the builder-confirmed mean take-off angle ψ, the "
        "per-detector effective angles follow as ψ ± s·θ "
        "(pending confirmation)"], size=22))
    set_notes(s, "BACKUP - not part of the 8:00. Use if the abstract-vs-talk "
                 "difference on viewing angles comes up: the abstract promised "
                 "per-detector angles; with one tilt only the response is "
                 "identifiable, and the per-detector values need the nominal "
                 "head geometry from the instrument builder.")

    # ---- drop the nine template slides ----------------------------------
    for t in T:
        delete_slide(prs, t)
    prs.save(out)
    print("saved %s (%d slides)" % (out, len(prs.slides)))
    return out


if __name__ == "__main__":
    ap = argparse.ArgumentParser(description=__doc__.split("\n")[2])
    ap.add_argument("--template", default=TEMPLATE)
    ap.add_argument("--out", default=OUT)
    args = ap.parse_args()
    build(args.template, args.out)
