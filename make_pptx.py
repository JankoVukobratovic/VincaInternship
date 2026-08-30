"""
Generates presentation.pptx for the XRF vulnerability assessment paper.
Run: python3 make_pptx.py
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import os

# ── paths ────────────────────────────────────────────────────────────────────
BASE = os.path.dirname(os.path.abspath(__file__))
FIG  = os.path.join(BASE, "xrf-denoise", "experiments", "full_pipeline", "figures")

IMG_ELEMENTS = os.path.join(FIG, "01_element_maps_raw_vs_denoised.png")
IMG_NMF      = os.path.join(FIG, "02_nmf_components.png")
IMG_CVI      = os.path.join(FIG, "03_cvi_map.png")
IMG_RULES    = os.path.join(FIG, "04_risk_rules.png")
IMG_SAM      = os.path.join(FIG, "05_sam_segmentation.png")

# ── palette ──────────────────────────────────────────────────────────────────
C_BG       = RGBColor(0x0D, 0x1B, 0x2A)   # deep navy
C_ACCENT   = RGBColor(0x4F, 0xC3, 0xF7)   # sky blue
C_GOLD     = RGBColor(0xF5, 0xC5, 0x18)   # amber
C_WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
C_LGRAY    = RGBColor(0xCC, 0xD6, 0xE0)
C_RED      = RGBColor(0xFF, 0x5A, 0x5A)
C_GREEN    = RGBColor(0x4C, 0xAF, 0x50)
C_ORANGE   = RGBColor(0xFF, 0x99, 0x33)
C_DARKBOX  = RGBColor(0x1A, 0x2D, 0x40)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width  = SLIDE_W
prs.slide_height = SLIDE_H

blank_layout = prs.slide_layouts[6]   # completely blank


# ── helpers ──────────────────────────────────────────────────────────────────

def bg(slide, color=C_BG):
    """Fill slide background."""
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def box(slide, left, top, width, height, text="", font_size=18,
        bold=False, color=C_WHITE, align=PP_ALIGN.LEFT,
        bg_color=None, bg_alpha=None):
    txb = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height))
    if bg_color:
        fill = txb.fill
        fill.solid()
        fill.fore_color.rgb = bg_color
    tf = txb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    return txb

def title_bar(slide, title, subtitle=""):
    """Dark accent bar at top of content slides."""
    bar = slide.shapes.add_shape(
        1, Inches(0), Inches(0), SLIDE_W, Inches(1.15))
    bar.fill.solid()
    bar.fill.fore_color.rgb = C_DARKBOX
    bar.line.fill.background()

    box(slide, 0.25, 0.07, 12.8, 0.65, title,
        font_size=30, bold=True, color=C_ACCENT, align=PP_ALIGN.LEFT)
    if subtitle:
        box(slide, 0.25, 0.72, 12.8, 0.38, subtitle,
            font_size=16, color=C_LGRAY, align=PP_ALIGN.LEFT)

def divider(slide, y, color=C_ACCENT, width=12.8, left=0.25):
    line = slide.shapes.add_shape(
        1, Inches(left), Inches(y), Inches(width), Inches(0.03))
    line.fill.solid()
    line.fill.fore_color.rgb = color
    line.line.fill.background()

def bullet_block(slide, left, top, width, height, items,
                 font_size=17, color=C_WHITE, indent="  • "):
    txb = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txb.text_frame
    tf.word_wrap = True
    first = True
    for item in items:
        if first:
            p = tf.paragraphs[0]; first = False
        else:
            p = tf.add_paragraph()
        p.space_before = Pt(4)
        run = p.add_run()
        run.text = indent + item
        run.font.size = Pt(font_size)
        run.font.color.rgb = color

def add_image(slide, path, left, top, width=None, height=None):
    if width and height:
        slide.shapes.add_picture(path, Inches(left), Inches(top),
                                 Inches(width), Inches(height))
    elif width:
        slide.shapes.add_picture(path, Inches(left), Inches(top),
                                 Inches(width))
    elif height:
        slide.shapes.add_picture(path, Inches(left), Inches(top),
                                 height=Inches(height))

def colored_label(slide, left, top, w, h, text, bg_c, font_size=13, text_color=C_WHITE):
    shp = slide.shapes.add_shape(1, Inches(left), Inches(top), Inches(w), Inches(h))
    shp.fill.solid(); shp.fill.fore_color.rgb = bg_c
    shp.line.fill.background()
    tf = shp.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    run = p.add_run(); run.text = text
    run.font.size = Pt(font_size); run.font.color.rgb = text_color; run.font.bold = True


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 - TITLE
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
bg(s)

# gradient bar left edge
bar = s.shapes.add_shape(1, Inches(0), Inches(0), Inches(0.18), SLIDE_H)
bar.fill.solid(); bar.fill.fore_color.rgb = C_ACCENT; bar.line.fill.background()

box(s, 0.5, 0.55, 12.3, 1.5,
    "Automated Chemical Vulnerability Assessment\nof Canvas Paintings from XRF Spectral Imaging\nUsing Deep Learning and Foundation Models",
    font_size=27, bold=True, color=C_WHITE)

divider(s, 2.25)

box(s, 0.5, 2.4, 12.3, 0.4,
    "Dimitrije Pešić  ·  Janko Vukobratović  ·  Aleksandra Stojanović  ·  Giulia Ristori  ·  Stefano Ridolfi  ·  Maja Gajić-Kvaščev  ·  Goran Kvaščev",
    font_size=13, color=C_LGRAY)

box(s, 0.5, 2.85, 12.3, 0.35,
    "University of Belgrade, School of Electrical Engineering  ·  Ars Mensurae Rome  ·  IDArtScience Rome  ·  Vinča Institute of Nuclear Sciences",
    font_size=12, color=C_LGRAY)

# key stats strip
stats = [("7,200", "XRF spectra"), ("< 70 s", "end-to-end"), ("r > 0.98", "reproducibility"), ("0 labels", "training data")]
for i, (val, lab) in enumerate(stats):
    xpos = 0.5 + i * 3.2
    shp = s.shapes.add_shape(1, Inches(xpos), Inches(5.6), Inches(2.9), Inches(1.4))
    shp.fill.solid(); shp.fill.fore_color.rgb = C_DARKBOX; shp.line.fill.background()
    box(s, xpos+0.08, 5.62, 2.75, 0.7, val, font_size=28, bold=True, color=C_ACCENT, align=PP_ALIGN.CENTER)
    box(s, xpos+0.08, 6.28, 2.75, 0.35, lab, font_size=13, color=C_LGRAY, align=PP_ALIGN.CENTER)

box(s, 0.5, 7.1, 12.3, 0.3, "Science Conference 2026", font_size=12, color=C_LGRAY)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 - MOTIVATION
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
bg(s)
title_bar(s, "Motivation", "Why automate conservation risk assessment?")

# Left column - problem
box(s, 0.3, 1.3, 4.2, 0.4, "THE PROBLEM", font_size=13, bold=True, color=C_ACCENT)
bullet_block(s, 0.3, 1.72, 4.3, 4.5, [
    "Historical paintings degrade silently - damage often invisible until irreversible",
    "Pigments interact chemically: Cu greens corrode, Pb whites turn black, Fe catalyzes oxidation",
    "Traditional assessment: expert-driven, qualitative, does not scale to full surfaces",
], font_size=16)

# Center column - XRF
box(s, 4.9, 1.3, 4.0, 0.4, "XRF MACRO-SCANNING", font_size=13, bold=True, color=C_ACCENT)
bullet_block(s, 4.9, 1.72, 4.0, 4.5, [
    "Per-pixel elemental maps across entire painting surface",
    "120 × 60 = 7,200 spectra at 3 s/pixel",
    "Non-invasive, no sampling required",
    "Problem: interpreting 7,200 spectra manually is impractical",
], font_size=16)

# Right column - our answer
box(s, 9.2, 1.3, 3.8, 0.4, "OUR ANSWER", font_size=13, bold=True, color=C_GOLD)
bullet_block(s, 9.2, 1.72, 3.9, 4.5, [
    "Fully automated pipeline",
    "Raw XRF → conservation-priority map",
    "Zero expert-labeled training data",
    "Under 70 s on a standard laptop",
], font_size=16, color=C_GOLD)

# vertical dividers
for xd in [4.75, 9.05]:
    dv = s.shapes.add_shape(1, Inches(xd), Inches(1.3), Inches(0.03), Inches(5.5))
    dv.fill.solid(); dv.fill.fore_color.rgb = C_DARKBOX; dv.line.fill.background()

divider(s, 7.1, C_DARKBOX)
box(s, 0.3, 7.15, 12.8, 0.3,
    "Key principle: preserve as much original material as possible → detect risk before it is visible",
    font_size=13, color=C_LGRAY, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 - PIPELINE OVERVIEW
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
bg(s)
title_bar(s, "Pipeline Overview", "Five stages · no labeled training data · <70 s end-to-end")

stages = [
    ("1", "Self-Supervised\nDenoising", "1D U-Net\nPoisson splitting", C_ACCENT),
    ("2", "Elemental\nExtraction", "5-pt calibration\nROI + background sub.", C_ACCENT),
    ("3", "NMF\nCorroboration", "K=5 components\n<5% recon. error", RGBColor(0x78,0xC8,0x78)),
    ("4", "Chemical\nVulnerability\nIndex", "5 literature rules\nGeometric mean", C_GOLD),
    ("5", "SAM Region\nSegmentation", "13 regions\nPer-region CVI", C_RED),
]

box_w = 2.1
gap   = 0.25
start = 0.35

for i, (num, title_t, sub, col) in enumerate(stages):
    x = start + i * (box_w + gap)
    # box
    shp = s.shapes.add_shape(1, Inches(x), Inches(1.5), Inches(box_w), Inches(3.8))
    shp.fill.solid(); shp.fill.fore_color.rgb = C_DARKBOX; shp.line.fill.background()
    # colored top strip
    top_strip = s.shapes.add_shape(1, Inches(x), Inches(1.5), Inches(box_w), Inches(0.22))
    top_strip.fill.solid(); top_strip.fill.fore_color.rgb = col; top_strip.line.fill.background()
    # number
    box(s, x, 1.72, box_w, 0.5, num, font_size=32, bold=True, color=col, align=PP_ALIGN.CENTER)
    # title
    box(s, x, 2.18, box_w, 0.85, title_t, font_size=15, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    # sub
    box(s, x+0.05, 3.1, box_w-0.1, 1.1, sub, font_size=12, color=C_LGRAY, align=PP_ALIGN.CENTER)
    # arrow (except last)
    if i < len(stages)-1:
        ax = x + box_w + 0.04
        box(s, ax, 3.05, gap+0.15, 0.5, "→", font_size=22, bold=True, color=C_LGRAY, align=PP_ALIGN.CENTER)

# inputs / outputs
box(s, 0.35, 5.55, 2.1, 0.35, "INPUT: raw .mca spectra",
    font_size=12, color=C_LGRAY, align=PP_ALIGN.CENTER)
box(s, start + 4*(box_w+gap), 5.55, box_w, 0.35, "OUTPUT: per-region risk report",
    font_size=12, color=C_GOLD, align=PP_ALIGN.CENTER)

# runtime row
divider(s, 6.15, C_DARKBOX)
runtime = [("U-Net inference", "28 s · 40%"), ("Extraction + NMF", "10.5 s · 15%"),
           ("CVI + metrics", "3.5 s · 5%"), ("SAM segmentation", "28 s · 40%")]
for i, (stage_n, t) in enumerate(runtime):
    x2 = 0.5 + i * 3.2
    box(s, x2, 6.25, 3.0, 0.3, stage_n, font_size=12, color=C_LGRAY, align=PP_ALIGN.CENTER)
    box(s, x2, 6.58, 3.0, 0.35, t, font_size=14, bold=True, color=C_ACCENT, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 - SPECTRAL ANALYSIS
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
bg(s)
title_bar(s, "Spectral Analysis", "Denoising → calibration → element maps → NMF corroboration")

# LEFT: denoising
box(s, 0.3, 1.25, 5.5, 0.38, "SELF-SUPERVISED DENOISING", font_size=13, bold=True, color=C_ACCENT)
bullet_block(s, 0.3, 1.65, 5.5, 2.5, [
    "XRF counts follow Poisson statistics",
    "Noise2Self: Binomial-split spectrum into two halves (x, y)",
    "Train 1D U-Net to predict y from x → no clean reference needed",
    "4-level encoder/decoder, ~1.66M params, PNLL loss",
], font_size=14)

box(s, 0.3, 4.0, 5.5, 0.38, "ELEMENTAL EXTRACTION", font_size=13, bold=True, color=C_ACCENT)
bullet_block(s, 0.3, 4.4, 5.5, 2.6, [
    "5-point linear calibration: E = 0.0292·c − 0.069 keV  (R² > 0.9999)",
    "ROI integration + linear background from off-peak sidebands",
    "Spectral overlap corrections (Cu Kβ→Zn, Pb/As regression)",
    "As/Pb Lα overlap at ΔE = 0.007 keV - below detector FWHM",
], font_size=14)

# RIGHT: element maps image
if os.path.exists(IMG_ELEMENTS):
    add_image(s, IMG_ELEMENTS, left=6.0, top=1.2, width=7.1)

# NMF footer
divider(s, 6.8, C_DARKBOX)
box(s, 0.3, 6.87, 12.7, 0.35,
    "NMF (K=5) independently recovers the same 5 material groups - Pb, Ca/Fe, Cu, Ti, Compton - reconstruction error <5%   →   cross-validates physics-based extraction",
    font_size=13, color=C_LGRAY, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 - CVI
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
bg(s)
title_bar(s, "Chemical Vulnerability Index (CVI)", "Literature-grounded degradation risk, pixel-level")

# Formula box
fbox = s.shapes.add_shape(1, Inches(0.3), Inches(1.25), Inches(5.6), Inches(0.85))
fbox.fill.solid(); fbox.fill.fore_color.rgb = C_DARKBOX; fbox.line.fill.background()
box(s, 0.35, 1.28, 5.5, 0.8,
    "CVI(i,j) = max_k { w_k · √( A_k(i,j) · B_k(i,j) ) }",
    font_size=17, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

bullet_block(s, 0.3, 2.25, 5.6, 2.0, [
    "Geometric mean → high score only when BOTH incompatible elements co-occur",
    "Single-element rules (R2, R3): CVI_k = w_k · A_k",
    "All intensities percentile-normalised → CVI in [0, 1]",
    "Gaussian filter σ=1 px to suppress pixel-level noise",
], font_size=14)

# Zones
box(s, 0.3, 4.35, 5.6, 0.35, "FOUR RISK ZONES", font_size=13, bold=True, color=C_ACCENT)
zones = [
    ("Low  <0.25",       C_GREEN,   "No action"),
    ("Moderate 0.25–0.50", RGBColor(0x8B,0xC3,0x4A), "Monitor"),
    ("Elevated 0.50–0.75", C_ORANGE, "Intervene"),
    ("Critical ≥0.75",   C_RED,     "Urgent"),
]
for i, (zname, zcol, zact) in enumerate(zones):
    colored_label(s, 0.3 + i*1.38, 4.75, 1.3, 0.55, zname, zcol, font_size=11)
    box(s, 0.3 + i*1.38, 5.35, 1.3, 0.3, zact, font_size=12, color=C_LGRAY, align=PP_ALIGN.CENTER)

# Rules table (right column)
box(s, 6.1, 1.25, 6.9, 0.38, "DEGRADATION RULES", font_size=13, bold=True, color=C_ACCENT)

rows = [
    ("R1", "TiO₂/CaCO₃ thermal mismatch",       "Ti / Ca", "0.40", C_LGRAY),
    ("R2", "Cu-based green pigment degradation", "Cu",      "1.00", C_RED),
    ("R3", "Lead white darkening (PbS)",         "Pb",      "1.00", C_RED),
    ("R4", "Trapped moisture under TiO₂",        "Ti / Cu", "0.60", C_ORANGE),
    ("R5", "Fe-catalyzed Pb oxidation",          "Fe / Pb", "0.80", C_ORANGE),
]

hdr_y = 1.72
for col_x, col_w, col_txt in [(6.1,0.5,"ID"),(6.65,4.0,"Mechanism"),(10.7,1.0,"Elem."),(11.75,0.9,"w")]:
    box(s, col_x, hdr_y, col_w, 0.35, col_txt, font_size=13, bold=True, color=C_ACCENT)

divider(s, 2.12, C_DARKBOX, width=6.9, left=6.1)

for i, (rid, mech, elem, w, col) in enumerate(rows):
    ry = 2.2 + i * 0.72
    if i % 2 == 0:
        rb = s.shapes.add_shape(1, Inches(6.1), Inches(ry-0.05), Inches(6.9), Inches(0.68))
        rb.fill.solid(); rb.fill.fore_color.rgb = C_DARKBOX; rb.line.fill.background()
    box(s, 6.1,  ry, 0.5, 0.55, rid,  font_size=14, bold=True, color=col)
    box(s, 6.65, ry, 4.0, 0.55, mech, font_size=13, color=C_WHITE)
    box(s, 10.7, ry, 1.0, 0.55, elem, font_size=13, color=C_LGRAY)
    box(s, 11.75,ry, 0.9, 0.55, w,    font_size=14, bold=True, color=col)

# sensitivity note
divider(s, 6.75, C_DARKBOX)
box(s, 0.3, 6.82, 12.7, 0.35,
    "Sensitivity: ±15% weight perturbation → <8% change in mean CVI  ·  top-decile ranking preserved in >92% of pixels",
    font_size=13, color=C_LGRAY, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 - RESULTS
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
bg(s)
title_bar(s, "Results", "Reproducibility · CVI distribution · SAM region analysis")

# ── left panel: reproducibility ──
box(s, 0.3, 1.25, 4.1, 0.38, "SCAN REPRODUCIBILITY", font_size=13, bold=True, color=C_ACCENT)
box(s, 0.3, 1.65, 4.1, 0.28, "Two scans, 7 days apart, same conditions", font_size=13, color=C_LGRAY)

repro = [("Ca Kα","0.992"),("Fe Kα","0.993"),("Pb Lα","0.991"),("Cu Kα","0.983"),("Ti Kα","0.982")]
for i, (el, r) in enumerate(repro):
    ry2 = 2.05 + i * 0.47
    rb2 = s.shapes.add_shape(1, Inches(0.3), Inches(ry2), Inches(4.1), Inches(0.43))
    rb2.fill.solid(); rb2.fill.fore_color.rgb = C_DARKBOX; rb2.line.fill.background()
    box(s, 0.38, ry2+0.05, 2.5, 0.35, el, font_size=14, color=C_WHITE)
    box(s, 2.95, ry2+0.05, 1.2, 0.35, r,  font_size=14, bold=True, color=C_GREEN)

# composite CVI metrics
box(s, 0.3, 4.45, 4.1, 0.35, "Composite CVI (prova1 vs prova2)", font_size=13, color=C_LGRAY)
for label, val, col2 in [("W₁","0.0054",C_GREEN),("r","0.994",C_GREEN),("SSIM","0.982",C_GREEN)]:
    pass  # will do as a row
metric_box = s.shapes.add_shape(1, Inches(0.3), Inches(4.82), Inches(4.1), Inches(0.55))
metric_box.fill.solid(); metric_box.fill.fore_color.rgb = C_DARKBOX; metric_box.line.fill.background()
box(s, 0.38, 4.87, 1.2, 0.4, "W₁=0.0054", font_size=13, bold=True, color=C_GREEN)
box(s, 1.65, 4.87, 1.1, 0.4, "r=0.994",   font_size=13, bold=True, color=C_GREEN)
box(s, 2.85, 4.87, 1.4, 0.4, "SSIM=0.982",font_size=13, bold=True, color=C_GREEN)

# ablation note
abox = s.shapes.add_shape(1, Inches(0.3), Inches(5.5), Inches(4.1), Inches(0.75))
abox.fill.solid(); abox.fill.fore_color.rgb = C_DARKBOX; abox.line.fill.background()
box(s, 0.38, 5.52, 3.9, 0.7,
    "Without denoising: W₁ +65%, SSIM −0.018\n→ U-Net improves stability but physics-based\n   extraction is dominant source of agreement",
    font_size=12, color=C_LGRAY)

# ── center: CVI zones ──
box(s, 4.75, 1.25, 3.8, 0.38, "CVI ZONE DISTRIBUTION", font_size=13, bold=True, color=C_ACCENT)
zone_data = [("Low",      "13.9%", C_GREEN,  1.65),
             ("Moderate", "49.9%", RGBColor(0x8B,0xC3,0x4A), 2.45),
             ("Elevated", "30.3%", C_ORANGE, 3.25),
             ("Critical", "5.9%",  C_RED,    4.05)]
for zn, zpct, zcol, zy in zone_data:
    zshp = s.shapes.add_shape(1, Inches(4.75), Inches(zy), Inches(3.8), Inches(0.72))
    zshp.fill.solid(); zshp.fill.fore_color.rgb = C_DARKBOX; zshp.line.fill.background()
    # color dot
    dot = s.shapes.add_shape(9, Inches(4.85), Inches(zy+0.2), Inches(0.32), Inches(0.32))
    dot.fill.solid(); dot.fill.fore_color.rgb = zcol; dot.line.fill.background()
    box(s, 5.28, zy+0.17, 2.0, 0.4, zn,   font_size=15, color=C_WHITE)
    box(s, 7.0,  zy+0.17, 1.3, 0.4, zpct, font_size=18, bold=True, color=zcol, align=PP_ALIGN.RIGHT)

box(s, 4.75, 4.9, 3.8, 0.65,
    "R2 (Cu degradation) elevated on ~22% of pixels\nR3 (Pb darkening) on ~19%\nR5 traces figural contours",
    font_size=13, color=C_LGRAY)

# CVI image
if os.path.exists(IMG_CVI):
    add_image(s, IMG_CVI, left=4.75, top=5.6, width=3.8)

# ── right: SAM regions ──
box(s, 8.85, 1.25, 4.2, 0.38, "SAM REGION ANALYSIS", font_size=13, bold=True, color=C_ACCENT)
box(s, 8.85, 1.65, 4.2, 0.28, "13 regions auto-identified", font_size=13, color=C_LGRAY)

sam_segs = [
    ("Pb-dominated", "max CVI 0.93 · mean 0.54", "Urgent (R3 – Pb darkening)", C_RED),
    ("Cu-dominated", "max CVI 0.97 · mean 0.45", "Monitor (R2 – Cu degr.)",    C_ORANGE),
    ("Ti/Cu overlap","mean CVI ≈ 0.45",            "Moisture entrapment (R4)",   RGBColor(0x8B,0xC3,0x4A)),
]
for i, (seg, metrics, risk, col3) in enumerate(sam_segs):
    sy = 2.05 + i * 1.35
    sshp = s.shapes.add_shape(1, Inches(8.85), Inches(sy), Inches(4.2), Inches(1.25))
    sshp.fill.solid(); sshp.fill.fore_color.rgb = C_DARKBOX; sshp.line.fill.background()
    left_bar = s.shapes.add_shape(1, Inches(8.85), Inches(sy), Inches(0.12), Inches(1.25))
    left_bar.fill.solid(); left_bar.fill.fore_color.rgb = col3; left_bar.line.fill.background()
    box(s, 9.05, sy+0.05, 3.9, 0.38, seg,     font_size=15, bold=True, color=C_WHITE)
    box(s, 9.05, sy+0.42, 3.9, 0.3,  metrics, font_size=13, color=C_LGRAY)
    box(s, 9.05, sy+0.72, 3.9, 0.35, risk,    font_size=13, color=col3)

if os.path.exists(IMG_SAM):
    add_image(s, IMG_SAM, left=8.85, top=6.1, width=4.2)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 - FURTHER WORK
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
bg(s)
title_bar(s, "Further Work & Conclusions", "Limitations and next steps")

# Summary column
box(s, 0.3, 1.25, 5.8, 0.38, "WHAT WE SHOWED", font_size=13, bold=True, color=C_GREEN)
bullet_block(s, 0.3, 1.65, 5.8, 3.5, [
    "Fully automated XRF → conservation-priority map in <70 s",
    "No expert-labeled training data at any stage",
    "Cross-scan reproducibility r > 0.98 on all elements",
    "Composite CVI stable across scans (W₁ = 0.0054, SSIM = 0.982)",
    "13 SAM regions with per-region risk summaries",
    "CVI critical zone: 5.9% of pixels flagged",
], font_size=16)

divider(s, 5.3, C_DARKBOX, width=5.8, left=0.3)
box(s, 0.3, 5.38, 5.8, 0.38, "PRINCIPAL LIMITATION", font_size=13, bold=True, color=C_ORANGE)
box(s, 0.3, 5.78, 5.8, 0.65,
    "All validation is internal consistency on one mockup canvas - ground-truth comparison against conservator assessments on real paintings with documented degradation outcomes is required.",
    font_size=14, color=C_LGRAY)

# Further work column
box(s, 7.0, 1.25, 5.9, 0.38, "NEXT STEPS", font_size=13, bold=True, color=C_ACCENT)

steps = [
    (C_RED,    "Validate CVI against expert conservator assessments on real paintings with known degradation outcomes"),
    (C_ORANGE, "Integrate NMF abundance maps as CVI inputs - capture material mixtures invisible to single-element thresholds"),
    (C_ACCENT, "Extend to paintings with different conservation histories; adapt rule weights as needed"),
    (C_ACCENT, "Resolve As/Pb Lα overlap (ΔE=0.007 keV) with higher-resolution detectors or deconvolution"),
    (RGBColor(0x8B,0xC3,0x4A), "Build conservator-facing GUI/API - no Python expertise required"),
]
for i, (col4, txt) in enumerate(steps):
    sy2 = 1.75 + i * 1.05
    sbox = s.shapes.add_shape(1, Inches(7.0), Inches(sy2), Inches(5.9), Inches(0.95))
    sbox.fill.solid(); sbox.fill.fore_color.rgb = C_DARKBOX; sbox.line.fill.background()
    lbar = s.shapes.add_shape(1, Inches(7.0), Inches(sy2), Inches(0.1), Inches(0.95))
    lbar.fill.solid(); lbar.fill.fore_color.rgb = col4; lbar.line.fill.background()
    box(s, 7.18, sy2+0.1, 5.6, 0.78, txt, font_size=13, color=C_WHITE)

divider(s, 7.0, C_ACCENT)
box(s, 0.3, 7.07, 12.7, 0.35,
    "The contribution is the workflow itself: end-to-end automation of a procedure that traditionally requires extensive expert analysis.",
    font_size=14, bold=True, color=C_ACCENT, align=PP_ALIGN.CENTER)


# ── save ─────────────────────────────────────────────────────────────────────
out = os.path.join(BASE, "presentation.pptx")
prs.save(out)
print(f"Saved: {out}")